"""
演示文稿解析器（PPTX / PPT / WPS DPS）
- PPTX: 使用 python-pptx 读取，提取每页幻灯片的文本和表格
- PPT/DPS: 旧版二进制格式，python-pptx 不支持，降级提示
每页幻灯片输出为一个 text block，保留页码信息
"""
import os
from loguru import logger
from typing import List, Dict, Any

from .base import BaseParser


class PresentationParser(BaseParser):
    name = 'presentation'

    def parse(self, file_path: str, **options) -> List[Dict[str, Any]]:
        ext = os.path.splitext(file_path)[1].lower()

        if ext in ('.pptx', '.dps'):
            return self._parse_pptx(file_path)
        elif ext == '.ppt':
            # .ppt 是旧版二进制格式，python-pptx 不支持
            logger.warning(f'[PresentationParser] .ppt 旧格式支持有限: {file_path}')
            return self._parse_ppt_fallback(file_path)
        else:
            return self._parse_pptx(file_path)

    def _parse_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        """使用 python-pptx 解析 PPTX/DPS 文件"""
        try:
            from pptx import Presentation
        except ImportError:
            logger.error('[PresentationParser] python-pptx 未安装，无法解析 PPTX')
            return []

        try:
            prs = Presentation(file_path)
        except Exception:
            logger.exception(f'[PresentationParser] 打开 PPTX 失败: {file_path}')
            return []

        blocks: List[Dict[str, Any]] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            texts = []
            tables = []

            for shape in slide.shapes:
                # 提取文本框
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)

                # 提取表格
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip().replace('|', '\\|') for cell in row.cells]
                        rows.append(' | '.join(cells))
                    if rows:
                        # 表头后加分隔行
                        if len(rows) > 0:
                            col_count = len(table.columns)
                            rows.insert(1, ' | '.join(['---'] * col_count))
                        tables.append('\n'.join(rows))

            # 幻灯片标题（取第一个文本作为标题）
            title = texts[0][:64] if texts else f'幻灯片 {slide_num}'

            # 文本内容
            if texts:
                content = '\n'.join(texts)
                blocks.append({
                    'type': 'text',
                    'content': content[:10000],
                    'section_path': f'幻灯片 {slide_num}: {title}',
                    'page_number': slide_num,
                    'extra': {'slide_number': slide_num, 'title': title},
                })

            # 表格内容
            for i, table_content in enumerate(tables):
                blocks.append({
                    'type': 'table',
                    'content': table_content[:20000],
                    'section_path': f'幻灯片 {slide_num}: {title} - 表格{i + 1}',
                    'page_number': slide_num,
                    'extra': {
                        'slide_number': slide_num,
                        'table_index': i,
                        'format': 'pptx_table',
                    },
                })

        return blocks

    def _parse_ppt_fallback(self, file_path: str) -> List[Dict[str, Any]]:
        """.ppt 旧版二进制格式降级处理"""
        logger.warning(f'[PresentationParser] .ppt 旧版格式无法解析，建议转换为 .pptx: {file_path}')
        return [{
            'type': 'text',
            'content': '[此文件为旧版 .ppt 格式，系统暂不支持自动解析。请将文件另存为 .pptx 格式后重新上传。]',
            'section_path': '解析提示',
            'page_number': None,
            'extra': {'format': 'ppt', 'parse_error': 'unsupported_legacy_format'},
        }]
