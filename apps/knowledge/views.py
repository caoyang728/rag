"""
knowledge views
- 节点树 & CRUD
- 文档上传（sha256 去重 -> 存盘 -> 触发 parse_document Celery 任务）
- 文档 chunks 查看
"""
import hashlib
import difflib
import magic
import os
import re
from loguru import logger
import uuid as uuid_lib

from django.conf import settings
from django.db import transaction, models
from django.db.models import Count
from django.http import FileResponse, Http404
from django.utils import text as django_text
from django.utils import timezone
from rest_framework import viewsets, status
from rest_framework.decorators import action
from rest_framework.exceptions import PermissionDenied, ValidationError
from rest_framework.parsers import MultiPartParser, FormParser, JSONParser
from rest_framework.permissions import IsAuthenticated
from rest_framework.response import Response
from rest_framework.views import APIView

from apps.knowledge.models import (
    KnowledgeNode, Document, DocumentChunk, DocOperationLog,
    ResourceShare, ResourceBlockList,
    VisibilityLevel, ResourceType, ShareScopeType, AccessLevel, InheritMode, ShareStatus,
)
from apps.knowledge.serializers import (
    KnowledgeNodeSerializer, KnowledgeNodeCreateSerializer,
    DocumentSerializer, DocumentChunkSerializer,
)
from apps.knowledge.access import resolve_doc_access, build_user_context, build_grants_map
from apps.knowledge.storage import get_document_storage
# 权限体系：ResourceShare/ResourceBlockList(文档级共享与黑名单) + TicketList(统一审批工单)
from apps.users.models import (
    User, Role, UserRoleRel, TicketList, TicketPermissionDetail, TicketFlowLog,
    TicketStatus, TicketChangeType, TicketBizType, ScopeType,
    has_permission, get_user_managed_teams, get_user_managed_depts,
)
from apps.users.permissions import IsAdminOrOps


def _log_operation(request, action, document=None, node=None, detail=None):
    """记录操作日志"""
    try:
        user = request.user if hasattr(request, 'user') and request.user.is_authenticated else None
        xff = request.META.get('HTTP_X_FORWARDED_FOR', '')
        ip = xff.split(',')[0].strip() if xff else request.META.get('REMOTE_ADDR', '')
        ua = request.META.get('HTTP_USER_AGENT', '')[:512]
        DocOperationLog.objects.create(
            action=action,
            operator=user,
            operator_name=user.username if user else '',
            document=document,
            node=node,
            detail=detail or {},
            ip_address=ip or None,
            user_agent=ua,
        )
    except Exception:
        logger.exception("log operation failed")


def _get_user_role(user):
    """获取用户角色信息（基于 RBAC permission_key 判定，清除角色硬编码）

    返回 (role, dept_id, team_ids)：
    - role: super_admin / kb_admin / dept_manager / team_leader / employee
    - dept_id: 用户主部门 ID
    - team_ids: 用户可管理的团队 ID 列表（含本团队 + 授权团队）

    判定优先级：super_admin → kb_admin → dept_manager(有 user.manage 且管理部门) → team_leader(有 user.manage 且管理团队) → employee
    """
    if not user or not getattr(user, 'is_authenticated', False):
        return None, None, []

    # super_admin 系统级快路径
    if getattr(user, 'is_super_admin', False):
        return 'super_admin', None, []

    # kb_admin：有 kb.manage_all 权限（is_kb_admin 属性内部已走 has_permission）
    if getattr(user, 'is_kb_admin', False):
        return 'kb_admin', None, []

    dept_id = getattr(user, 'department_id', None)
    # user.team 为单团队 FK，managed_teams 含本团队 + UserTeamScopeRel 授权团队
    managed_team_ids = list(get_user_managed_teams(user))
    managed_dept_ids = get_user_managed_depts(user)

    # 部门管理员：有 user.manage 权限且可管理部门多于主部门（即有属地授权）
    if has_permission(user, 'user.manage') and managed_dept_ids:
        return 'dept_manager', dept_id, managed_team_ids

    # 团队组长：有 user.manage 权限且管理团队非空
    if has_permission(user, 'user.manage') and managed_team_ids:
        return 'team_leader', dept_id, managed_team_ids

    # 只有显式授权 contributor 的用户才能获得写权限分类
    # viewer 兜底用户返回 None，调用方会将其视为无上传权限
    has_contributor_role = Role.objects.filter(
        role_key='contributor',
        id__in=UserRoleRel.objects.filter(
            user=user, status='ACTIVE',
        ).values_list('role_id', flat=True),
    ).exists()
    if has_contributor_role:
        return 'contributor', dept_id, managed_team_ids
    return None, dept_id, managed_team_ids


# 旧 visible_scope → 新 visibility_level 映射（兼容前端旧参数）
_LEGACY_SCOPE_MAP = {
    'team': VisibilityLevel.TEAM_ONLY,
    'dept': VisibilityLevel.DEPT_ONLY,
    'public': VisibilityLevel.PUBLIC,
}


def _normalize_visibility_level(value):
    """将前端传入的可见性参数归一化为 VisibilityLevel 枚举值

    兼容两种输入：
    - 新版 visibility_level: TEAM_ONLY / DEPT_ONLY / PUBLIC
    - 旧版 visible_scope: team / dept / public
    返回 VisibilityLevel 枚举值或 None（非法值）
    """
    if not value:
        return None
    if value in VisibilityLevel.values:
        return value
    return _LEGACY_SCOPE_MAP.get(value)


def _resolve_node_visibility(node):
    """从节点沿祖先链向上取最近非空可见范围（NULL=继承父级）

    用于文档上传未显式指定可见范围时继承挂载文件夹的可见性，
    保证"节点可见范围是文档可见性的默认值来源"这一语义：
    - 文件夹设置了可见范围 → 取其值
    - 文件夹未设置 → 继续向上取部门/团队节点
    - 全部未设置 → root 兜底 PUBLIC
    """
    cur = node
    seen = set()
    while cur and cur.id not in seen:
        seen.add(cur.id)
        if cur.visibility_level:
            return cur.visibility_level
        cur = cur.parent
    return VisibilityLevel.PUBLIC


def _resolve_node_org(node):
    """沿祖先链解析节点所属部门/团队的组织 ID（非节点 ID）

    节点树分层：root(1) → 部门节点(2) → 团队节点(3) → 文件夹(4+)。
    部门/团队节点由组织同步创建（node_kind=ORG 且带 ref_id 指向组织表记录），
    因此沿祖先链向上匹配 node_kind=ORG 的节点即可还原节点归属。

    返回 (dept_id, team_id)：
    - 挂载在部门节点下的文件夹 → (dept_id, None)
    - 挂载在团队节点下的文件夹 → (dept_id, team_id)
    - root 下的纯公共文件夹（无 ORG 祖先）→ (None, None)

    用途：可见范围变更工单按发起人角色动态指派审批链时，
    需要知道目标节点的组织归属以定位审批人（部门经理等）。
    """
    cur = node
    seen = set()
    dept_id = None
    team_id = None
    while cur and cur.id not in seen:
        seen.add(cur.id)
        if cur.node_kind == 'ORG':
            if cur.node_level == 2 and cur.ref_id:
                dept_id = cur.ref_id
            elif cur.node_level == 3 and cur.ref_id:
                team_id = cur.ref_id
        cur = cur.parent
    return dept_id, team_id


def _build_visibility_chain(applicant, node):
    """按发起人角色 + 节点组织归属动态构建节点可见范围变更审批链

    对齐 plan.md T1 验收"审批链按角色指派正确"：
    - 团队级：团队组长发起（节点在其团队子树内）→ 部门经理审批
    - 部门级：部门经理发起（节点在其部门子树内）→ 文档管理员/超管审批
    - 超管/知识库管理员发起 → 双管理员复核（保留既有双层审批语义）

    审批链节点结构兼容现有 approval_chain 消费逻辑
    （step/approver_id/status/comment/approved_at），
    并带 approver_role / approver_scope_id 供审批时按角色匹配审批人。
    """
    def _step(role, scope_id=None):
        """构造单个审批链节点（step 序号由调用方回填）"""
        return {
            'step': 0,
            'approver_role': role,
            'approver_scope_id': scope_id,
            'approver_id': None,
            'status': 'pending',
            'comment': '',
            'approved_at': None,
        }

    role, _dept_id, _team_ids = _get_user_role(applicant)
    node_dept_id, _node_team_id = _resolve_node_org(node)

    if role == 'team_leader' and node_dept_id:
        # 团队级：部门经理审批（scope 锁定节点所属部门）
        chain = [_step('DEPT_LEADER', node_dept_id)]
    elif role == 'dept_manager' and node_dept_id:
        # 部门级：文档管理员/超管审批
        chain = [_step('KB_ADMIN')]
    else:
        # 超管/知识库管理员发起或组织归属缺失：双管理员复核（保留既有语义）
        chain = [_step('KB_ADMIN'), _step('KB_ADMIN')]
    for i, n in enumerate(chain):
        n['step'] = i
    return chain


def _get_dept_node_paths(user) -> list:
    """用户有属地授权（管理/上传）的部门节点 path 列表（Level 2 ORG）

    通过 get_user_managed_depts 获取部门 ID 后反查部门节点，
    用于"节点是否在本部门子树内"的 path 前缀匹配判定（部门经理领地）。
    权限判定与上传/写文件夹共用，避免各处重复查询。
    """
    from apps.users.models import get_user_managed_depts
    dept_ids = get_user_managed_depts(user)
    if not dept_ids:
        return []
    return list(KnowledgeNode.objects.filter(
        node_level=2, ref_id__in=dept_ids, is_deleted=False,
    ).values_list('path', flat=True))


def _get_team_node_paths(team_ids) -> list:
    """给定团队 ID 列表的团队节点 path 列表（Level 3 ORG）

    用于"节点是否在本团队子树内"的 path 前缀匹配判定（团队组长领地）。
    """
    if not team_ids:
        return []
    return list(KnowledgeNode.objects.filter(
        node_level=3, ref_id__in=team_ids, is_deleted=False,
    ).values_list('path', flat=True))


def _can_approve_node_visibility(user, node, current_step):
    """判断用户是否有权审批节点可见范围工单的当前审批步骤

    动态审批链（approver_role 非空）时按角色匹配审批人：
    - DEPT_LEADER：节点所属部门的部门经理（user.manage + 属地授权含该部门），
      超管/文档管理员兜底可审
    - KB_ADMIN：文档管理员（kb.manage_all）或超管
    - SUPER_ADMIN：仅超管
    旧工单（无 approver_role，兼容固定管理员链）：
    节点所有者 / 超管 / 文档管理员均可审。
    """
    approver_role = (current_step or {}).get('approver_role')
    is_admin = bool(getattr(user, 'is_super_admin', False)
                    or getattr(user, 'is_kb_admin', False))
    if not approver_role:
        return (node.owner_user_id == user.id or is_admin)
    if approver_role == 'DEPT_LEADER':
        if is_admin:
            return True
        dept_id = (current_step or {}).get('approver_scope_id')
        if not dept_id or not has_permission(user, 'user.manage'):
            return False
        return dept_id in get_user_managed_depts(user)
    if approver_role == 'KB_ADMIN':
        return is_admin
    if approver_role == 'SUPER_ADMIN':
        return bool(getattr(user, 'is_super_admin', False))
    return is_admin


def _validate_visibility_level(user, visibility_level):
    """验证用户是否有权限设置指定的可见性层级

    返回 (is_valid, error_message)

    visibility_level 三档:
    - TEAM_ONLY: 仅归属团队可见
    - DEPT_ONLY: 仅归属部门（含下属团队）可见
    - PUBLIC: 全局全员可见（扩大范围需审批）
    """
    if visibility_level not in VisibilityLevel.values:
        return False, "无效的可见性层级设置"

    # super_admin / kb_admin 可以设置任意可见性层级
    if getattr(user, 'is_super_admin', False) or getattr(user, 'is_kb_admin', False):
        return True, None

    # 所有用户都可以设置 TEAM_ONLY 和 DEPT_ONLY
    if visibility_level in (VisibilityLevel.TEAM_ONLY, VisibilityLevel.DEPT_ONLY):
        return True, None

    # PUBLIC 需要审批（通过统一工单流程），但创建时可以设置
    return True, None


def _encode_ticket_reason(target_type, target_id, action, user_reason=''):
    """将访问申请/可见范围变更的目标信息编码到工单 reason 字段

    权限工单详情（TicketPermissionDetail）不区分文档/节点目标，
    文档/节点访问申请复用统一工单，通过 reason 前缀编码目标。
    格式: [doc:{doc_id}:{action}] 或 [node:{node_id}:{action}] 后接用户申请理由
    """
    return f"[{target_type}:{target_id}:{action}] {user_reason or ''}"


def _decode_ticket_reason(reason):
    """从工单 reason 解析目标信息

    返回 (target_type, target_id, action, user_reason)：
    - target_type: 'doc' 文档工单 / 'node' 节点工单 / None 无法解析
    无法解析时返回 (None, None, None, reason)
    """
    if not reason:
        return None, None, None, ''
    import re as _re
    m = _re.match(r'^\[(doc|node):(\d+):(\w+)\]\s*(.*)$', reason, _re.DOTALL)
    if m:
        return m.group(1), int(m.group(2)), m.group(3), m.group(4)
    return None, None, None, reason


def _extract_last_comment(approval_chain):
    """从工单 approval_chain 提取最近一条审批意见

    approval_chain 格式: [{step, approver_id, status, comment, approved_at}, ...]
    用于在申请列表中展示审批人意见。
    """
    if not approval_chain or not isinstance(approval_chain, list):
        return ''
    # 从后往前找第一条有 comment 的记录
    for step in reversed(approval_chain):
        if isinstance(step, dict) and step.get('comment'):
            return step['comment']
    return ''


def _create_doc_ticket(applicant, target_user, change_type, reason, chain):
    """创建文档访问/可见性变更统一工单 —— 主表 + 权限详情子表 + SUBMIT 流转日志

    文档/节点目标信息编码在 reason（见 _encode_ticket_reason），
    详情子表记录发起人/目标人/变更类型，role 为空（文档域工单不涉及角色授权）。
    主表走统一工单号（QX + 日期 + 当日序列），保证与工单中心一致。
    """
    from apps.users.ticket_service import _gen_ticket_no
    ticket = TicketList.objects.create(
        ticket_no=_gen_ticket_no(TicketBizType.PERMISSION),
        title=f'文档权限·{change_type}'.strip(),
        biz_type=TicketBizType.PERMISSION,
        status=TicketStatus.PENDING,
        risk_level='normal',
        applicant=applicant,
        approval_chain=chain,
        current_step=0,
    )
    TicketPermissionDetail.objects.create(
        ticket=ticket,
        target_user=target_user,
        change_type=change_type,
        role=None,
        scope_type=ScopeType.NONE,
        scope_id=None,
        reason=reason,
    )
    TicketFlowLog.objects.create(ticket=ticket, action='SUBMIT', actor=applicant)
    return ticket


ALLOWED_EXTENSIONS = {
    # 文档类
    ".pdf", ".doc", ".docx", ".md", ".markdown", ".txt", ".rst",
    # 电子表格
    ".csv", ".xlsx", ".xls",
    # 演示文稿
    ".ppt", ".pptx",
    # 代码类
    ".py", ".java", ".go", ".js", ".ts", ".jsx", ".tsx", ".c", ".cpp", ".h", ".rs",
    # 配置类
    ".yaml", ".yml", ".json", ".xml", ".toml", ".ini", ".conf", ".cfg",
    # 脚本/样式
    ".sh", ".bat", ".ps1", ".css",
    # WPS Office 格式（新版WPS默认使用标准格式，这些是兼容旧版扩展名）
    ".wps", ".et", ".dps",
}
MAX_FILE_SIZE = int(getattr(settings, 'DOCUMENT_MAX_SIZE_MB', 100)) * 1024 * 1024

FILE_TYPE_MAP = {
    ".pdf": "pdf",
    ".doc": "docx", ".docx": "docx",
    ".wps": "docx",  # WPS 文字，尝试用 docx 解析
    ".md": "markdown", ".markdown": "markdown",
    ".txt": "txt",
    ".csv": "spreadsheet", ".xlsx": "spreadsheet", ".xls": "spreadsheet",
    ".et": "spreadsheet",  # WPS 表格，尝试用 spreadsheet 解析
    ".ppt": "presentation", ".pptx": "presentation",
    ".dps": "presentation",  # WPS 演示，尝试用 presentation 解析
    ".py": "code", ".java": "code", ".go": "code", ".js": "code",
    ".ts": "code", ".c": "code", ".cpp": "code", ".rs": "code",
    ".yaml": "config", ".yml": "config", ".json": "config",
    ".toml": "config", ".ini": "config", ".conf": "config",
}


def _detect_file_type(filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower()
    return FILE_TYPE_MAP.get(ext, "other")


# ============================================================================
# 文档活跃版本判定（同组文件：「新版本」 vs 「恰好同名的独立文档」）
# ============================================================================
# 文本类文件可直接读取内容做相似度判定；二进制（pdf/docx/xlsx 等）上传时无法
# 即时提取文本，默认按「新版本」处理（见 _is_version_upload）。
_VERSION_TEXT_FILE_TYPES = ('txt', 'markdown', 'code', 'config')
# 同组文件内容相似度 >= 对应阈值 → 视为同一文档的新版本（旧版本自动置非活跃）；
# 低于阈值 → 视为恰好同名的独立文档（全部保留活跃，如不同项目的同名代码文件）。
# 阈值按文件类型区分：
# - txt/markdown 业务文档：多为增量更新（如年会名单），内容相近即视为新版本
# - code/config 代码配置：同名文件常为不同项目的独立实现（如两个项目的 views.py），
#   需内容高度相似（近乎同一文件）才视为新版本，避免互相覆盖
VERSION_SIMILARITY_THRESHOLD = 0.3
_VERSION_SIMILARITY_THRESHOLDS = {
    'txt': 0.3,
    'markdown': 0.3,
    'code': 0.8,
    'config': 0.8,
}
# 内容样本截取上限：先按原始字节截断（避免解码超大文件），再按字符截断。
VERSION_SAMPLE_MAX_BYTES = 8192
VERSION_SAMPLE_MAX_CHARS = 4000


def _capture_content_sample(raw_bytes, file_type):
    """从上传文件原始字节中截取规范化文本样本，用于版本相似度判定

    仅文本类文件（txt/markdown/code/config）可即时解码；二进制文件返回空串，
    调用方会将其按「新版本」处理。样本统一折叠空白，换行/缩进差异不影响相似度。
    """
    if file_type not in _VERSION_TEXT_FILE_TYPES or not raw_bytes:
        return ''
    text = raw_bytes.decode('utf-8', errors='ignore')
    return re.sub(r'\s+', ' ', text).strip()[:VERSION_SAMPLE_MAX_CHARS]


def _text_similarity(s1, s2):
    """两个文本样本的相似度 [0,1]（difflib 序列比对）

    用于区分「同一文档的新版本」（相似度高）与「恰好同名的独立文档」（相似度低）。
    """
    if not s1 or not s2:
        return 0.0
    return difflib.SequenceMatcher(None, s1, s2).ratio()


def _is_version_upload(file_type, content_sample, siblings):
    """判断本次上传是否为同组已有文档的「新版本」

    返回 True 时，同组（node+file_name+dept_id+team_id）其他文档将自动置非活跃。
    - 无同组文档：False（首传，无需置非活跃）
    - 二进制文件 / 无文本样本：True（无法即时读文，保守按新版本处理）
    - 文本类：与同组任一文档样本相似度 >= 对应文件类型阈值 → True，否则 False（视为独立文档）
    """
    if not siblings:
        return False
    if file_type not in _VERSION_TEXT_FILE_TYPES or not content_sample:
        return True
    threshold = _VERSION_SIMILARITY_THRESHOLDS.get(file_type, VERSION_SIMILARITY_THRESHOLD)
    return any(
        sib.content_sample
        and _text_similarity(content_sample, sib.content_sample) >= threshold
        for sib in siblings
    )


def _sync_vectors_active(doc_ids, is_active):
    """同步文档的检索向量活跃标志（DocumentVector 冗余字段，检索层无需 JOIN 主表）

    在活跃版本切换（上传新版本 / set_active）后调用，保证检索层即时生效。
    """
    if not doc_ids:
        return
    from apps.retrieval.models import DocumentVector
    DocumentVector.objects.filter(document_id__in=list(doc_ids)).update(is_active=is_active)


def _extract_text_content(content: bytes, file_type: str, filename: str) -> str:
    """从文件内容中提取文本（用于预览）"""
    ext = os.path.splitext(filename)[1].lower()

    # 文本类文件：直接解码
    if file_type in ("txt", "markdown") or ext in (".txt", ".md", ".markdown"):
        try:
            return content.decode("utf-8")
        except UnicodeDecodeError:
            try:
                return content.decode("gbk")
            except UnicodeDecodeError:
                return content.decode("latin-1", errors="ignore")

    # PDF
    if file_type == "pdf" or ext == ".pdf":
        try:
            from pypdf import PdfReader
            import io
            reader = PdfReader(io.BytesIO(content))
            text = "\n\n".join(page.extract_text() or "" for page in reader.pages)
            return text if text.strip() else "PDF 文档无文本内容"
        except ImportError:
            return "需要安装 pypdf 才能预览 PDF 内容"
        except Exception as e:
            logger.error(f"PDF extract failed: {e}")
            return f"PDF 解析失败: {str(e)}"

    # DOCX
    if file_type == "docx" or ext in (".doc", ".docx"):
        try:
            from docx import Document
            import io
            doc = Document(io.BytesIO(content))
            text = "\n\n".join(paragraph.text for paragraph in doc.paragraphs)
            return text if text.strip() else "Word 文档无文本内容"
        except ImportError:
            return "需要安装 python-docx 才能预览 Word 内容"
        except Exception as e:
            logger.error(f"DOCX extract failed: {e}")
            return f"Word 解析失败: {str(e)}"

    # 代码/配置文件
    if file_type in ("code", "config"):
        try:
            return content.decode("utf-8")
        except UnicodeDecodeError:
            try:
                return content.decode("gbk")
            except UnicodeDecodeError:
                return content.decode("latin-1", errors="ignore")

    # 电子表格（CSV/XLSX/XLS/ET）
    if file_type == "spreadsheet" or ext in (".csv", ".xlsx", ".xls", ".et"):
        return _extract_spreadsheet_preview(content, ext)

    # 演示文稿（PPTX/PPT/DPS）
    if file_type == "presentation" or ext in (".ppt", ".pptx", ".dps"):
        return _extract_presentation_preview(content, ext)

    # 未知类型：尝试文本解码或显示二进制提示
    try:
        return content.decode("utf-8")
    except UnicodeDecodeError:
        return f"[{filename}] 无法预览此类型文件，建议下载查看"


def _extract_spreadsheet_preview(content: bytes, ext: str) -> str:
    """提取电子表格预览文本"""
    import io

    # CSV: 直接解码为文本
    if ext == '.csv':
        for encoding in ('utf-8', 'gbk', 'gb2312', 'latin-1'):
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue
        return "CSV 文件解码失败"

    # XLSX/ET: 使用 openpyxl 读取
    if ext in ('.xlsx', '.et'):
        try:
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                parts.append(f"=== Sheet: {ws.title} ===")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else '' for c in row]
                    if any(c.strip() for c in cells):
                        parts.append(' | '.join(cells))
            wb.close()
            return '\n'.join(parts) if parts else "电子表格无数据"
        except ImportError:
            return "需要安装 openpyxl 才能预览 Excel 内容"
        except Exception as e:
            logger.error(f"Spreadsheet preview failed: {e}")
            return f"电子表格解析失败: {str(e)}"

    # XLS: 旧版格式不支持预览
    if ext == '.xls':
        return "[.xls 旧版格式不支持预览，请下载查看或转换为 .xlsx 格式]"

    return "不支持的电子表格格式"


def _extract_presentation_preview(content: bytes, ext: str) -> str:
    """提取演示文稿预览文本"""
    import io

    # PPTX/DPS: 使用 python-pptx 读取
    if ext in ('.pptx', '.dps'):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            parts = []
            for slide_num, slide in enumerate(prs.slides, start=1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                texts.append(text)
                if texts:
                    parts.append(f"=== 幻灯片 {slide_num} ===")
                    parts.append('\n'.join(texts))
            return '\n'.join(parts) if parts else "演示文稿无文本内容"
        except ImportError:
            return "需要安装 python-pptx 才能预览 PPT 内容"
        except Exception as e:
            logger.error(f"Presentation preview failed: {e}")
            return f"演示文稿解析失败: {str(e)}"

    # PPT: 旧版格式不支持预览
    if ext == '.ppt':
        return "[.ppt 旧版格式不支持预览，请下载查看或转换为 .pptx 格式]"

    return "不支持的演示文稿格式"


def _build_tree(qs):
    """扁平列表 -> 递归树"""
    nodes = list(qs)
    id2n = {n["id"]: {**n, "children": []} for n in nodes}
    roots = []
    for n in nodes:
        pid = n.get("parent_id")
        if pid and pid in id2n:
            id2n[pid]["children"].append(id2n[n["id"]])
        else:
            roots.append(id2n[n["id"]])
    return roots


class NodeTreeView(APIView):
    """GET /api/v1/knowledge/nodes/tree/?root_type=company_doc"""
    permission_classes = [IsAuthenticated]

    def get(self, request):
        root_type = request.query_params.get("root_type")
        qs = KnowledgeNode.objects.filter(is_deleted=False)
        if root_type:
            qs = qs.filter(root_type=root_type)
        # 统计每个节点下未删除的文档数（与详情页 document_count 口径一致）
        qs = qs.annotate(
            document_count=Count(
                "documents",
                filter=models.Q(documents__is_deleted=False),
            ),
        )
        nodes = qs.order_by("depth", "order_no", "id").values(
            "id", "parent_id", "root_type", "node_type", "node_kind", "name", "depth",
            "node_level", "document_count", "ref_id", "path",
        )
        data = list(nodes)
        return Response({"tree": _build_tree(data), "total": len(data)})


class RootTypesView(APIView):
    """GET /api/v1/knowledge/nodes/root_types/ - 动态获取所有根类型"""
    permission_classes = [IsAuthenticated]

    def get(self, request):
        root_types = KnowledgeNode.objects.filter(
            node_type='root', is_deleted=False
        ).values_list('root_type', flat=True).distinct()
        return Response({
            "root_types": [
                {"code": rt, "name": rt} for rt in root_types
            ]
        })


class AllowedVisibilityView(APIView):
    """GET /api/v1/knowledge/documents/allowed_visibility/ - 获取当前用户可选的部门/团队"""
    permission_classes = [IsAuthenticated]

    def get(self, request):
        from apps.users.models import Department, Team
        from django.core.cache import cache

        role, user_dept_id, user_team_ids = _get_user_role(request.user)
        
        # 构建缓存key：不同角色有不同的缓存
        cache_key = f'allowed_visibility_{role}_{user_dept_id}_{tuple(sorted(user_team_ids))}'
        
        # 尝试从缓存获取
        cached_result = cache.get(cache_key)
        if cached_result is not None:
            return Response(cached_result)

        result = {
            'role': role,
            'can_set_public': True,  # 所有角色都可以设置公开
            'departments': [],
            'teams': [],
        }

        result['departments'] = list(Department.objects.filter(is_deleted=False).values('id', 'name'))
        result['teams'] = list(Team.objects.filter(is_deleted=False).values('id', 'name', 'code', 'department_id'))

        # 缓存1小时（3600秒）
        cache.set(cache_key, result, 3600)
        
        return Response(result)


class KnowledgeNodeViewSet(viewsets.ModelViewSet):
    """/api/v1/knowledge/nodes/ — retrieve 全员可见；写操作 admin/ops 或团队组长（仅本团队范围内）"""
    queryset = KnowledgeNode.objects.filter(is_deleted=False).order_by("path")
    serializer_class = KnowledgeNodeSerializer
    filterset_fields = ["root_type", "node_type", "parent"]

    # ── 团队组长权限辅助方法 ────────────

    def _get_team_leader_paths(self, user):
        """获取团队组长管理的团队节点 path 列表

        user.team 为单团队 FK，组长可能管理多个团队（通过 Team.leader 反查）。
        通过 ref_id=team.id 定位 Level 3 团队节点，取其 path 用于子树匹配。
        """
        try:
            from apps.users.models import Team
            team_ids = list(Team.objects.filter(
                leader=user, is_deleted=False
            ).values_list('id', flat=True))
        except Exception:
            team_ids = []
        return _get_team_node_paths(team_ids)

    def _check_team_node_write(self, node, user):
        """检查团队组长是否有权操作该节点（节点必须在组长团队子树内）"""
        team_paths = self._get_team_leader_paths(user)
        for tp in team_paths:
            # tp 本身以 / 结尾；子节点 path 以 tp 开头即为团队子树内
            if node.path == tp or node.path.startswith(tp):
                return
        raise PermissionDenied("您只能操作自己团队范围内的分类节点")

    def _is_admin_user(self, user):
        """用户是否为管理员（RBAC：knowledge:manage:all）"""
        try:
            return bool(user.is_kb_admin)
        except Exception:
            return False

    def _check_dept_node_write(self, node, user):
        """部门经理只能在本部门节点及其文件夹后代下创建/操作文件夹

        允许挂载：本部门节点本身（node_level=2 的 ORG，与部门节点同级挂文件夹）或
        本部门节点下的任意文件夹（path 前缀匹配）；
        拒绝：团队节点（组长领地）、其他部门节点、root 及越界文件夹。
        """
        if node.node_kind == 'ORG':
            # ORG 节点中只有本部门节点本身允许挂载文件夹（团队节点归组长管理）
            from apps.users.models import get_user_managed_depts
            if node.node_level == 2 and node.ref_id in get_user_managed_depts(user):
                return
            raise PermissionDenied("您只能在自己部门下创建文件夹")
        for dp in _get_dept_node_paths(user):
            if node.path == dp or node.path.startswith(dp):
                return
        raise PermissionDenied("您只能在自己部门下创建文件夹")

    # ── 权限 ────────────

    def get_permissions(self):
        """retrieve 面向所有登录用户开放；写操作允许 admin/ops 或团队组长"""
        if self.action == 'retrieve':
            return [IsAuthenticated()]
        if self.action in ('create', 'destroy', 'update', 'partial_update'):
            # 需要登录，具体权限在方法内校验
            return [IsAuthenticated()]
        return [IsAdminOrOps()]

    # ── 节点保护：ROOT/ORG 节点禁止通过节点 API 直接 CRUD ────────────
    # 根节点与组织节点（部门/团队）由系统自动创建，只有手动创建的文件夹（FOLDER）可被 CRUD
    _MANAGED_LABELS = {'ROOT': '根节点', 'ORG': '组织节点'}

    def _check_node_not_managed(self, node):
        """ROOT/ORG 节点由系统/组织架构同步管理，禁止通过节点 API 直接增删改"""
        if node.node_kind and node.node_kind != 'FOLDER':
            label = self._MANAGED_LABELS.get(node.node_kind, '系统节点')
            raise ValidationError(
                {'detail': f'{label}不支持直接操作，请通过部门/团队管理功能操作'}
            )

    def get_queryset(self):
        qs = super().get_queryset()
        # 预计算 children_count 和 document_count，避免 N+1
        if self.action in ("retrieve", "list"):
            qs = qs.annotate(
                _children_count=Count("children", filter=models.Q(children__is_deleted=False)),
                _document_count=Count("documents", filter=models.Q(documents__is_deleted=False)),
            )
        return qs

    def get_serializer_class(self):
        if self.action in ("create", "update", "partial_update"):
            return KnowledgeNodeCreateSerializer
        return KnowledgeNodeSerializer

    def perform_create(self, serializer):
        user = self.request.user
        parent = serializer.validated_data.get("parent")
        node_type = serializer.validated_data.get("node_type", "folder")

        # root 节点由系统自动创建（node_sync.get_or_create_kb_root），禁止手动创建
        if node_type == "root":
            raise ValidationError({"parent": "根节点由系统自动创建，不支持手动创建"})

        # 节点（组织分支）只能由部门/团队同步创建；手动创建的一律是文件夹
        if not parent:
            raise ValidationError({"parent": "文件夹必须指定上级节点"})

        role, _dept_id, _team_ids = _get_user_role(user)

        if role in ('super_admin', 'kb_admin'):
            # 超管/文档管理员：可在任意位置创建文件夹（含 root 下，与部门节点同级）
            pass
        elif role == 'dept_manager':
            # 部门经理：只能在本部门节点及其文件夹后代下创建
            self._check_dept_node_write(parent, user)
        elif role == 'team_leader':
            # 团队组长：只能在本团队节点及其文件夹后代下创建
            self._check_team_node_write(parent, user)
        else:
            raise PermissionDenied("您没有创建文件夹的权限")

        root_type = parent.root_type
        depth = parent.depth + 1
        # node_level 按挂载位置动态计算：root 下=2（与部门节点同级）/部门下=3/团队下=4/嵌套依次递增
        obj = serializer.save(depth=depth, node_level=parent.node_level + 1,
                              node_type='folder', node_kind='FOLDER',
                              root_type=root_type, created_by=user)
        # 更新 path（ID 零填充 4 位，确保按数值顺序排序）
        padded_id = f"{obj.id:04d}"
        obj.path = f"{parent.path}{padded_id}/"
        obj.save(update_fields=["path"])
        _log_operation(self.request, 'node_create', node=obj,
                       detail={'name': obj.name, 'node_type': obj.node_type,
                               'node_kind': obj.node_kind, 'root_type': obj.root_type})

    def destroy(self, request, *args, **kwargs):
        node = self.get_object()

        # 非管理员：只能删除自己领地内的文件夹
        if not self._is_admin_user(request.user):
            role, _dept_id, _team_ids = _get_user_role(request.user)
            if role == 'dept_manager':
                self._check_dept_node_write(node, request.user)
            elif role == 'team_leader':
                self._check_team_node_write(node, request.user)
            else:
                raise PermissionDenied("您没有删除文件夹的权限")

        # ROOT/ORG 保护：禁止直接删除（由部门/团队生命周期管理）
        if node.node_kind and node.node_kind != 'FOLDER':
            label = self._MANAGED_LABELS.get(node.node_kind, '系统节点')
            return Response(
                {"detail": f"{label}不支持直接删除，请通过部门/团队管理功能操作"},
                status=status.HTTP_400_BAD_REQUEST,
            )

        # 检查是否存在子节点/文件夹
        child_count = KnowledgeNode.objects.filter(
            parent=node, is_deleted=False
        ).count()
        if child_count > 0:
            return Response(
                {"detail": f"该节点下存在 {child_count} 个子分类，请先删除所有子分类后再删除此节点"},
                status=status.HTTP_400_BAD_REQUEST,
            )

        # 递归检查该节点及其所有子孙节点下是否存在文档
        from apps.knowledge.node_sync import count_docs_in_subtree
        doc_count = count_docs_in_subtree(node.id)
        if doc_count > 0:
            return Response(
                {"detail": f"该分类下存在 {doc_count} 个文档，请先迁移或删除所有文档后再删除此节点"},
                status=status.HTTP_400_BAD_REQUEST,
            )

        # 软删除
        node.is_deleted = True
        node.save(update_fields=["is_deleted"])
        _log_operation(request, 'node_delete', node=node,
                       detail={'name': node.name, 'node_type': node.node_type,
                               'node_kind': node.node_kind})
        return Response(status=204)

    def perform_update(self, serializer):
        old_obj = self.get_object()

        # ROOT/ORG 节点禁止修改（由部门/团队生命周期管理）
        self._check_node_not_managed(old_obj)

        # 非管理员：只能编辑自己领地内的文件夹
        if not self._is_admin_user(self.request.user):
            role, _dept_id, _team_ids = _get_user_role(self.request.user)
            if role == 'dept_manager':
                self._check_dept_node_write(old_obj, self.request.user)
            elif role == 'team_leader':
                self._check_team_node_write(old_obj, self.request.user)
            else:
                raise PermissionDenied("您没有修改文件夹的权限")

        # 可见范围变更（visibility_level）必须走工单审批，不能直接修改
        new_visibility_level = serializer.validated_data.get(
            'visibility_level', old_obj.visibility_level
        )
        if new_visibility_level != old_obj.visibility_level:
            user = self.request.user
            if new_visibility_level is not None and new_visibility_level not in VisibilityLevel.values:
                raise ValidationError({'visibility_level': '可见范围必须是 TEAM_ONLY/DEPT_ONLY/PUBLIC 之一'})
            # 目标可见范围编码在 reason 的"目标值:"标记中，审批通过后由 approve_access_request 解析写回
            old_desc = dict(VisibilityLevel.choices).get(old_obj.visibility_level, '继承父级')
            new_desc = dict(VisibilityLevel.choices).get(
                new_visibility_level, '继承父级') if new_visibility_level else '继承父级'
            _create_doc_ticket(
                user, user, TicketChangeType.SCOPE_CHANGE,
                _encode_ticket_reason(
                    'node', old_obj.id, 'visibility_change',
                    f"申请将节点可见范围从「{old_desc}」调整为「{new_desc}」 "
                    f"目标值:{new_visibility_level or 'INHERIT'}"
                ),
                # 审批链按发起人角色动态指派：团队级=部门经理，部门级=文档管理员/超管，
                # 超管/文档管理员=双管理员复核（_build_visibility_chain）
                _build_visibility_chain(user, old_obj),
            )
            raise PermissionDenied(
                "修改节点可见范围需要审批，已自动提交审批工单，审批通过后可见范围生效"
            )

        old_data = {
            'name': old_obj.name,
            'description': old_obj.description,
            'order_no': old_obj.order_no,
        }
        new_obj = serializer.save()
        new_data = {
            'name': new_obj.name,
            'description': new_obj.description,
            'order_no': new_obj.order_no,
        }
        _log_operation(self.request, 'node_update', node=new_obj,
                       detail={'old': old_data, 'new': new_data})


class DocumentViewSet(viewsets.ModelViewSet):
    """/api/v1/knowledge/documents/"""
    queryset = Document.objects.order_by("-created_at")
    serializer_class = DocumentSerializer
    permission_classes = [IsAuthenticated]
    # 支持按 visibility_level / dept_id / team_id 过滤
    filterset_fields = ["node", "status", "file_type", "visibility_level",
                        "root_type", "owner", "is_deleted", "dept_id", "team_id"]
    search_fields = ["title", "file_name", "owner__username", "owner__real_name"]

    def get_queryset(self):
        # restored_by 用于序列化器 restored_by_name，一并 select_related 避免列表页 N+1
        qs = super().get_queryset().select_related("owner", "node", "restored_by")
        user = self.request.user
        include_deleted = self.request.query_params.get("include_deleted") == "true"

        if not include_deleted:
            qs = qs.filter(is_deleted=False)

        # 活跃版本过滤：列表默认只返回活跃版本（检索/浏览视角），?version=all 返回全部。
        # 仅作用于 list 动作——detail 动作（retrieve/versions/set_active/update/destroy）必须
        # 能访问非活跃版本（版本切换、回溯旧版本），越权由 get_object 的 can_read 拦截。
        if (self.action == 'list'
                and self.request.query_params.get("version") != "all"):
            qs = qs.filter(is_active=True)

        # 部门筛选：Document.dept_id 直接存储部门组织 ID（非节点 ID），无需再查节点
        dept_id = self.request.query_params.get("dept_id")
        if dept_id:
            qs = qs.filter(dept_id=dept_id)

        # 可见性筛选：前端传 visible_scope 或 visibility（team/dept/public），归一化为 visibility_level
        vis_raw = self.request.query_params.get("visible_scope") or self.request.query_params.get("visibility")
        if vis_raw:
            vis_level = _normalize_visibility_level(vis_raw)
            if vis_level:
                qs = qs.filter(visibility_level=vis_level)

        if self.request.query_params.get("discover"):
            # 发现模式：返回全部文档用于浏览与申请权限，
            # 但绝密(secret_level=4)文档的条目名仅 owner 和管理员可见
            if not (getattr(user, 'is_super_admin', False) or getattr(user, 'is_kb_admin', False)):
                qs = qs.exclude(
                    models.Q(secret_level=4) & ~models.Q(owner=user)
                )
            return qs
        if getattr(user, 'is_super_admin', False) or getattr(user, 'is_kb_admin', False):
            return qs
        # 非管理员：仅看自己上传的 + 全局公开文档（visibility_level=PUBLIC）
        # 跨团队/跨部门的文档需通过 ResourceShare 共享才能可见，由检索层 access.py 过滤
        qs = qs.filter(
            models.Q(owner=user) |
            models.Q(visibility_level=VisibilityLevel.PUBLIC)
        )
        return qs

    def get_serializer_context(self):
        ctx = super().get_serializer_context()
        ctx["_user_ctx"] = build_user_context(self.request.user)
        return ctx

    def list(self, request, *args, **kwargs):
        queryset = self.filter_queryset(self.get_queryset())
        page = self.paginate_queryset(queryset)
        if page is not None:
            ctx = self.get_serializer_context()
            ctx["_grants_map"] = build_grants_map(request.user, [d.id for d in page])
            ctx["_version_count_map"] = self._build_version_count_map(page)
            serializer = self.get_serializer(page, many=True, context=ctx)
            return self.get_paginated_response(serializer.data)
        serializer = self.get_serializer(queryset, many=True)
        return Response(serializer.data)

    def _build_version_count_map(self, docs):
        """单次查询统计每个文档同组（node+file_name+dept_id+team_id）的版本总数

        列表页「版本切换」入口需要知道某文档是否还有其他版本；同组判定与上传/去重
        逻辑一致（含非活跃版本、跨分页），一次 SQL 完成，避免逐行 N+1 查询。
        """
        ids = [d.id for d in docs if d.id]
        if not ids:
            return {}
        group_keys = list(
            Document.objects.filter(pk__in=ids, is_deleted=False)
            .values_list('node_id', 'file_name', 'dept_id', 'team_id')
        )
        if not group_keys:
            return {}
        q = models.Q()
        for node_id, file_name, dept_id, team_id in group_keys:
            q |= models.Q(node_id=node_id, file_name=file_name,
                          dept_id=dept_id, team_id=team_id)
        rows = (Document.objects.filter(is_deleted=False).filter(q)
                .values('node_id', 'file_name', 'dept_id', 'team_id')
                .annotate(cnt=models.Count('id')))
        counts = {(r['node_id'], r['file_name'], r['dept_id'], r['team_id']): r['cnt']
                  for r in rows}
        return {d.id: counts.get((d.node_id, d.file_name, d.dept_id, d.team_id), 1)
                for d in docs}

    @action(detail=True, methods=["get"])
    def versions(self, request, pk=None):
        """GET /documents/{id}/versions/ 获取文档的同组版本列表（版本切换弹窗数据源）

        返回同组（node+file_name+dept_id+team_id）全部未删除版本，含版本号、活跃标记、
        处理状态与是否可切换（is_owner）；不校验写权限，读取元信息即可。
        """
        doc = self.get_object()
        siblings = (Document.objects.filter(
            node=doc.node, file_name=doc.file_name,
            dept_id=doc.dept_id, team_id=doc.team_id, is_deleted=False,
        ).order_by('-version', '-created_at'))
        is_admin = getattr(request.user, 'is_super_admin', False) or getattr(request.user, 'is_kb_admin', False)
        data = [{
            'id': d.id,
            'title': d.title,
            'version': d.version,
            'version_tag': d.version_tag or '',
            'is_active': d.is_active,
            'status': d.status,
            'file_size': d.file_size,
            'created_at': d.created_at,
            # 仅 Owner / 管理员可执行「设为活跃」
            'is_owner': d.owner_id == request.user.id or is_admin,
        } for d in siblings]
        return Response({'documents': data})

    @action(detail=True, methods=["post"])
    def set_active(self, request, pk=None):
        """POST /documents/{id}/set_active/ 将指定版本设为活跃版本

        同组（node+file_name+dept_id+team_id）其他版本自动置非活跃，并同步检索向量表的
        活跃标志；仅文档 Owner 或管理员可操作。幂等：已是活跃版本直接返回成功。
        """
        doc = self.get_object()
        self._require_write(doc)
        if doc.is_deleted:
            return Response({"detail": "已删除文档不能设置为活跃版本"}, status=400)
        if not doc.is_active:
            with transaction.atomic():
                deactivated_qs = Document.objects.filter(
                    node=doc.node, file_name=doc.file_name,
                    dept_id=doc.dept_id, team_id=doc.team_id, is_deleted=False,
                ).exclude(id=doc.id)
                deactivated_ids = list(deactivated_qs.values_list('id', flat=True))
                deactivated_qs.update(is_active=False)
                _sync_vectors_active(deactivated_ids, False)
                doc.is_active = True
                doc.save(update_fields=['is_active', 'updated_at'])
                _sync_vectors_active([doc.id], True)
            _log_operation(request, 'doc_set_active', document=doc, node=doc.node,
                           detail={'doc_id': doc.id, 'version_tag': doc.version_tag,
                                   'deactivated_ids': deactivated_ids})
        return Response({'id': doc.id, 'is_active': True})

    @action(detail=False, methods=["get"])
    def available_depts(self, request):
        """获取部门列表（用于筛选），使用Redis缓存"""
        from django.core.cache import cache
        from apps.users.models import Department
        
        cache_key = "available_depts_list"
        cached_depts = cache.get(cache_key)
        
        if cached_depts is not None:
            return Response(cached_depts)
        
        depts = Department.objects.filter(is_deleted=False).values('id', 'name')
        dept_list = list(depts)
        
        cache.set(cache_key, dept_list, 3600)
        
        return Response(dept_list)

    def get_object(self):
        obj = super().get_object()
        # 读取级校验：至少需要 can_read 才能获取单条详情
        if not self._access(obj)["can_read"]:
            raise PermissionDenied("无权限查看此文档")
        return obj

    def _access(self, doc, user=None):
        return resolve_doc_access(user or self.request.user, doc,
                                  ctx=build_user_context(user or self.request.user))

    def _require_write(self, doc):
        """写操作（编辑/删除/分享/管理授权）仅限所有者或管理员"""
        a = self._access(doc)
        if not (a["is_owner"] or a["is_manager"]):
            raise PermissionDenied("仅文档所有者或管理员可执行此操作")

    def perform_update(self, serializer):
        old_obj = self.get_object()
        access = self._access(old_obj)

        # 检查是否在修改可见性层级（visibility_level 控制可见范围）
        new_visibility_level = serializer.validated_data.get(
            'visibility_level', old_obj.visibility_level
        )
        is_changing_visibility = new_visibility_level != old_obj.visibility_level

        if is_changing_visibility:
            user = self.request.user

            # 验证可见性层级是否合法
            is_valid, error_msg = _validate_visibility_level(user, new_visibility_level)
            if not is_valid:
                raise PermissionDenied(error_msg)

            # 可见性扩大（TEAM_ONLY→DEPT_ONLY / TEAM_ONLY→PUBLIC / DEPT_ONLY→PUBLIC）需要双层审批
            # scope_order 值越大可见范围越广
            scope_order = {
                VisibilityLevel.TEAM_ONLY: 0,
                VisibilityLevel.DEPT_ONLY: 1,
                VisibilityLevel.PUBLIC: 2,
            }
            if scope_order.get(new_visibility_level, 0) > scope_order.get(old_obj.visibility_level, 0):
                # 创建统一工单（文档目标信息编码在 reason 字段前缀，approval_chain 记录双层审批步骤）
                _create_doc_ticket(
                    user, user, TicketChangeType.SCOPE_CHANGE,
                    _encode_ticket_reason(
                        'doc', old_obj.id, 'visibility_change',
                        f"申请将文档可见性从「{old_obj.get_visibility_level_display()}」"
                        f"扩大为「{dict(VisibilityLevel.choices).get(new_visibility_level, new_visibility_level)}」"
                    ),
                    # 双层审批链：两位不同管理员先后审批
                    [
                        {'step': 0, 'approver_id': None, 'status': 'pending', 'comment': '', 'approved_at': None},
                        {'step': 1, 'approver_id': None, 'status': 'pending', 'comment': '', 'approved_at': None},
                    ],
                )
                raise PermissionDenied(
                    "扩大可见性层级需要双层审批，已自动提交审批工单，需两位管理员先后审批"
                )
        else:
            # 其他字段修改：需要写权限
            self._require_write(old_obj)

        old_data = {
            'visibility_level': old_obj.visibility_level,
            'allow_download': old_obj.allow_download,
            'allow_share': old_obj.allow_share,
            'title': old_obj.title,
            'node_id': old_obj.node_id,
        }
        new_obj = serializer.save()
        new_data = {
            'visibility_level': new_obj.visibility_level,
            'allow_download': new_obj.allow_download,
            'allow_share': new_obj.allow_share,
            'title': new_obj.title,
            'node_id': new_obj.node_id,
        }
        # 检测可见性变更
        if is_changing_visibility:
            _log_operation(self.request, 'doc_visibility_change', document=new_obj,
                           detail={'old': old_data, 'new': new_data})

    def destroy(self, request, *args, **kwargs):
        doc = self.get_object()
        self._require_write(doc)
        doc.is_deleted = True
        doc.delete_time = timezone.now()
        doc.save(update_fields=["is_deleted", "delete_time"])
        _log_operation(request, 'doc_delete', document=doc,
                       detail={'title': doc.title, 'file_name': doc.file_name})
        try:
            from apps.retrieval.vector_store import delete_by_document
            delete_by_document(doc.id)
        except Exception:
            logger.exception("delete vector failed")
        # 文档删除联动：清理该文档产生的图谱数据（失败不阻断删除流程）
        try:
            from apps.graph.sync import on_document_deleted
            on_document_deleted(doc.id)
        except Exception:
            logger.exception("graph sync on document delete failed")
        return Response(status=204)

    @action(detail=True, methods=["post"])
    def restore(self, request, pk=None):
        """恢复已删除的文档"""
        doc = self.get_object()
        self._require_write(doc)
        if not doc.is_deleted:
            return Response({"detail": "文档未被删除"}, status=400)
        
        doc.is_deleted = False
        doc.restored_at = timezone.now()
        doc.restored_by = request.user
        doc.save(update_fields=["is_deleted", "restored_at", "restored_by"])
        _log_operation(request, 'doc_restore', document=doc,
                       detail={'title': doc.title, 'file_name': doc.file_name})
        return Response({"ok": True})

    @action(detail=True, methods=["post"])
    def hard_delete(self, request, pk=None):
        """
        物理删除已删除的文档（删除物理文件）
        
        限制条件：
        - 文档必须已被逻辑删除（is_deleted=True）
        - 删除时间超过30天（DEBUG模式下不受限制）
        - 超过180天的已删除文档会被自动清理任务删除
        
        物理删除后无法恢复。
        """
        from django.conf import settings
        from apps.knowledge.storage import get_document_storage
        
        doc = self.get_object()
        self._require_write(doc)
        
        if not doc.is_deleted:
            return Response({"detail": "文档未被逻辑删除，请先执行逻辑删除"}, status=400)
        
        if not doc.file_path:
            return Response({"detail": "文档没有物理文件可删除"}, status=400)
        
        min_retention_days = 30
        if not settings.DEBUG and doc.delete_time:
            days_since_delete = (timezone.now() - doc.delete_time).days
            if days_since_delete < min_retention_days:
                remaining_days = min_retention_days - days_since_delete
                return Response({
                    "detail": f"文档删除不足 {min_retention_days} 天，还需等待 {remaining_days} 天才能物理删除",
                    "remaining_days": remaining_days,
                    "days_since_delete": days_since_delete
                }, status=403)
        
        storage = get_document_storage()
        try:
            storage.delete(doc.file_path)
            doc.file_path = ''
            doc.save(update_fields=['file_path'])
            _log_operation(request, 'doc_hard_delete', document=doc,
                           detail={'title': doc.title, 'file_name': doc.file_name})
            return Response({"ok": True})
        except Exception as e:
            logger.exception(f"Failed to hard delete file for doc={doc.id}")
            return Response({"detail": f"物理删除失败: {str(e)[:200]}"}, status=500)

    @action(detail=True, methods=["post"])
    def reparse(self, request, pk=None):
        """重新解析文档：删除旧向量/切片/代码块/图片资源，基于原文件重新解析"""
        doc = self.get_object()
        self._require_write(doc)
        doc.status = "pending"
        doc.error_message = ""
        doc.save(update_fields=["status", "error_message"])
        _log_operation(request, 'doc_reparse', document=doc,
                       detail={'title': doc.title})
        try:
            from apps.knowledge.tasks import parse_document
            parse_document.delay(doc.id)
        except Exception:
            logger.exception("dispatch parse task failed")
        return Response({"ok": True, "status": "pending"})

    # ------------------------------------------------------------------
    # 下载
    # ------------------------------------------------------------------
    @action(detail=True, methods=["get"])
    def download(self, request, pk=None):
        doc = self.get_object()
        if not self._access(doc)["can_download"]:
            raise PermissionDenied("无权限下载此文档")
        if not doc.file_path:
            raise Http404("文件不存在")
        _log_operation(request, 'doc_download', document=doc,
                       detail={'file_name': doc.file_name, 'file_size': doc.file_size})
        # OSS：返回签名 URL 跳转；本地：直接返回文件流
        if doc.file_path.startswith("oss://"):
            storage = get_document_storage()
            url = storage.get_url(doc.file_path)
            from django.http import HttpResponseRedirect
            return HttpResponseRedirect(url)
        if not os.path.exists(doc.file_path):
            raise Http404("文件不存在")
        fp = open(doc.file_path, "rb")
        return FileResponse(fp, as_attachment=True, filename=doc.file_name)

    # ------------------------------------------------------------------
    # 原始内容预览：返回文档原始文本内容（支持分页，用于前端预览，不可复制）
    # ------------------------------------------------------------------
    @action(detail=True, methods=["get"])
    def raw_content(self, request, pk=None):
        doc = self.get_object()
        if not self._access(doc)["can_read"]:
            raise PermissionDenied("无权限预览此文档")
        if not doc.file_path:
            raise Http404("文件不存在")

        # 分页参数（每页字符数）
        page = max(1, int(request.query_params.get("page", 1)))
        page_size = min(max(1000, int(request.query_params.get("page_size", 5000))), 20000)  # 1k-20k 字符

        # 获取完整文本内容（最大 50MB，超出则截断）
        MAX_PREVIEW_SIZE = 50 * 1024 * 1024
        text_content = self._get_document_text(doc, MAX_PREVIEW_SIZE)
        
        if text_content is None:
            return Response({"error": "无法获取文件内容"}, status=500)
        
        total_chars = len(text_content)
        total_pages = max(1, (total_chars + page_size - 1) // page_size)
        page = min(page, total_pages)
        
        # 计算当前页内容（按字符数分页，尽量在段落边界处断开）
        start = (page - 1) * page_size
        end = start + page_size
        
        # 尽量在换行符处断开
        if end < total_chars:
            # 向后找最近的换行符
            newline_pos = text_content.find('\n', end - 100, end + 200)
            if newline_pos != -1:
                end = newline_pos + 1  # 包含换行符
        
        current_content = text_content[start:end]
        
        # 添加上下文提示
        if start > 0:
            current_content = '...' + current_content
        if end < total_chars:
            current_content = current_content + '...'

        return Response({
            "content": current_content,
            "file_type": doc.file_type,
            "file_name": doc.file_name,
            "size": doc.file_size or total_chars,
            "total_chars": total_chars,
            "total_pages": total_pages,
            "current_page": page,
            "page_size": page_size,
            "can_copy": False,
        })

    def _get_document_text(self, doc, max_size):
        """提取文档文本内容（内部方法）"""
        content = None
        if doc.file_path.startswith("oss://"):
            storage = get_document_storage()
            url = storage.get_url(doc.file_path)
            try:
                import requests
                resp = requests.get(url, timeout=30, stream=True)
                content = b""
                for chunk in resp.iter_content(chunk_size=1024 * 1024):
                    content += chunk
                    if len(content) > max_size:
                        content = content[:max_size]
                        break
            except Exception as e:
                logger.error(f"OSS raw content fetch failed: {e}")
                return None
        else:
            if not os.path.exists(doc.file_path):
                raise Http404("文件不存在")
            with open(doc.file_path, "rb") as f:
                content = f.read(max_size)
        
        return _extract_text_content(content, doc.file_type, doc.file_name)

    # ------------------------------------------------------------------
    # 申请权限
    # ------------------------------------------------------------------
    @action(detail=True, methods=["post"])
    def request_access(self, request, pk=None):
        """POST /documents/{id}/request_access/  {action, reason?}
        注意：申请者通常尚无读取权限，故不走 get_object 的 can_read 校验。

        通过统一工单（TicketList）管理审批工单：
        文档目标信息编码在 reason 前缀 [doc:{doc_id}:{action}]。
        """
        doc = Document.objects.filter(id=pk, is_deleted=False).first()
        if not doc:
            raise Http404("文档不存在")
        action = request.data.get("action", "read")
        if action not in ("read", "download"):
            raise ValidationError({"action": "无效的申请类型"})

        user_reason = (request.data.get("reason") or "")[:1000]
        encoded_reason = _encode_ticket_reason('doc', doc.id, action, user_reason)

        # 已有相同 pending 申请则不重复创建（通过 reason 前缀匹配文档目标）
        exists = TicketList.objects.filter(
            applicant=request.user,
            status=TicketStatus.PENDING,
            permission_detail__change_type=TicketChangeType.GRANT,
            permission_detail__reason=encoded_reason,
        ).exists()
        if exists:
            return Response({"ok": False, "detail": "已存在待审批的相同申请"}, status=200)

        ticket = _create_doc_ticket(
            request.user, request.user, TicketChangeType.GRANT, encoded_reason,
            [
                {'step': 0, 'approver_id': None, 'status': 'pending', 'comment': '', 'approved_at': None},
            ],
        )
        logger.info(f"[AccessRequest] doc={doc.id} applicant={request.user.username} action={action}")
        return Response({
            "id": ticket.id,
            "doc_id": doc.id,
            "action": action,
            "reason": user_reason,
            "status": ticket.status,
            "created_at": ticket.created_at,
        }, status=201)

    # ------------------------------------------------------------------
    # 访问授权管理（所有者/管理员查看与撤销）
    # ------------------------------------------------------------------
    @action(detail=True, methods=["get"])
    def access_grants(self, request, pk=None):
        """GET /documents/{id}/access_grants/  查看该文档的所有授权

        统一查询 ResourceShare（共享白名单）+ ResourceBlockList（黑名单）：
        - allow_users: ResourceShare(share_scope_type=USER)
        - cross_teams: ResourceShare(share_scope_type=TEAM)
        - deny_users: ResourceBlockList（仅个人）
        """
        doc = self.get_object()
        self._require_write(doc)

        result = {
            'allow_users': [],      # 个人共享白名单
            'cross_teams': [],      # 跨团队共享
            'deny_users': [],       # 黑名单
            'visibility_level': doc.visibility_level,
        }

        # 1. ResourceShare（共享白名单：个人 + 团队）
        shares = ResourceShare.objects.filter(
            resource_type=ResourceType.DOCUMENT,
            resource_id=doc.id,
        ).select_related('granted_by')
        for share in shares:
            if share.share_scope_type == ShareScopeType.USER:
                result['allow_users'].append({
                    'id': share.id,
                    'uid': share.share_scope_id,
                    'access_level': share.access_level,
                    'status': share.status,
                    'expire_time': share.expires_at,
                    'create_time': share.granted_at,
                })
            elif share.share_scope_type == ShareScopeType.TEAM:
                # 查团队 code 用于展示
                from apps.users.models import Team
                team_code = Team.objects.filter(id=share.share_scope_id).values_list('code', flat=True).first()
                result['cross_teams'].append({
                    'id': share.id,
                    'team_id': share.share_scope_id,
                    'team_code': team_code or '',
                    'access_level': share.access_level,
                    'status': share.status,
                    'expire_time': share.expires_at,
                    'create_time': share.granted_at,
                })

        # 2. ResourceBlockList（黑名单，仅个人）
        for block in ResourceBlockList.objects.filter(
            resource_type=ResourceType.DOCUMENT, resource_id=doc.id
        ).select_related('blocked_user'):
            result['deny_users'].append({
                'id': block.id,
                'uid': block.blocked_user_id,
                'reason': block.reason,
                'status': block.status,
                'create_time': block.blocked_at,
            })

        return Response(result)

    @action(detail=True, methods=["post"], url_path="grant_access")
    def grant_access(self, request, pk=None):
        """POST /documents/{id}/grant_access/  {grant_type, team_code/uid}  创建跨团队/个人授权

        通过 ResourceShare 管理跨团队/个人授权：
        - cross_team: ResourceShare(share_scope_type=TEAM, share_scope_id=team_id)
        - allow_user: ResourceShare(share_scope_type=USER, share_scope_id=user_id)
        授予后设置 doc.has_resource_share=True 加速检索过滤。
        """
        doc = self.get_object()
        self._require_write(doc)
        grant_type = request.data.get("grant_type")
        if grant_type not in ("cross_team", "allow_user"):
            raise ValidationError({"grant_type": "无效的授权类型，可选: cross_team/allow_user"})
        try:
            if grant_type == "cross_team":
                team_code = request.data.get("team_code", "").strip()
                if not team_code:
                    raise ValidationError({"team_code": "team_code 不能为空"})
                # 通过 team_code 查 team_id
                from apps.users.models import Team
                team = Team.objects.filter(code=team_code, is_deleted=False).first()
                if not team:
                    raise ValidationError({"team_code": f"团队 {team_code} 不存在"})
                # ResourceShare 唯一约束：(resource_type, resource_id, share_scope_type, share_scope_id)
                # 撤销后重新授予会产生新记录，此处用 get_or_create 避免重复
                share, created = ResourceShare.objects.get_or_create(
                    resource_type=ResourceType.DOCUMENT,
                    resource_id=doc.id,
                    share_scope_type=ShareScopeType.TEAM,
                    share_scope_id=team.id,
                    defaults={
                        'access_level': AccessLevel.READ,
                        'granted_by': request.user,
                        'status': ShareStatus.ACTIVE,
                    },
                )
                if created and not doc.has_resource_share:
                    doc.has_resource_share = True
                    doc.save(update_fields=["has_resource_share"])
                _log_operation(request, 'doc_share', document=doc,
                               detail={'grant_type': 'cross_team', 'team_code': team_code,
                                       'team_id': team.id, 'created': created})
                return Response({
                    "id": share.id,
                    "grant_type": "cross_team",
                    "team_code": team_code,
                    "created": created,
                })
            else:
                uid = request.data.get("uid")
                if not uid:
                    raise ValidationError({"uid": "uid 不能为空"})
                share, created = ResourceShare.objects.get_or_create(
                    resource_type=ResourceType.DOCUMENT,
                    resource_id=doc.id,
                    share_scope_type=ShareScopeType.USER,
                    share_scope_id=uid,
                    defaults={
                        'access_level': AccessLevel.READ,
                        'granted_by': request.user,
                        'status': ShareStatus.ACTIVE,
                    },
                )
                if created and not doc.has_resource_share:
                    doc.has_resource_share = True
                    doc.save(update_fields=["has_resource_share"])
                _log_operation(request, 'doc_share', document=doc,
                               detail={'grant_type': 'allow_user', 'uid': uid,
                                       'created': created})
                return Response({
                    "id": share.id,
                    "grant_type": "allow_user",
                    "uid": uid,
                    "created": created,
                })
        except Exception as e:
            raise ValidationError({"detail": str(e)})

    @action(detail=True, methods=["post"], url_path="revoke_grant")
    def revoke_grant(self, request, pk=None):
        """POST /documents/{id}/revoke_grant/  {grant_type, grant_id}

        统一操作 ResourceShare / ResourceBlockList：
        - allow_user / cross_team: ResourceShare.status=REVOKED（软撤销，保留历史审计）
        - deny_user: ResourceBlockList.status=REVOKED
        """
        doc = self.get_object()
        self._require_write(doc)
        grant_type = request.data.get("grant_type")
        grant_id = request.data.get("grant_id")
        if grant_type not in ("allow_user", "cross_team", "deny_user"):
            raise ValidationError({"grant_type": "无效的授权类型，可选: allow_user/cross_team/deny_user"})
        try:
            if grant_type in ("allow_user", "cross_team"):
                # 共享白名单撤销（ResourceShare）
                grant = ResourceShare.objects.get(
                    id=grant_id,
                    resource_type=ResourceType.DOCUMENT,
                    resource_id=doc.id,
                )
                grant.status = ShareStatus.REVOKED
                grant.revoked_by = request.user
                grant.revoked_at = timezone.now()
                grant.save(update_fields=['status', 'revoked_by', 'revoked_at'])
                revoked_detail = {
                    'grant_type': grant_type,
                    'grant_id': grant.id,
                    'share_scope_type': grant.share_scope_type,
                    'share_scope_id': grant.share_scope_id,
                }
            else:
                # 黑名单撤销（ResourceBlockList）
                grant = ResourceBlockList.objects.get(
                    id=grant_id,
                    resource_type=ResourceType.DOCUMENT,
                    resource_id=doc.id,
                )
                grant.status = ShareStatus.REVOKED
                grant.revoked_by = request.user
                grant.revoked_at = timezone.now()
                grant.save(update_fields=['status', 'revoked_by', 'revoked_at'])
                revoked_detail = {
                    'grant_type': 'deny_user',
                    'grant_id': grant.id,
                    'blocked_user_id': grant.blocked_user_id,
                }
        except (ResourceShare.DoesNotExist, ResourceBlockList.DoesNotExist):
            raise Http404("授权记录不存在")
        _log_operation(request, 'doc_revoke', document=doc,
                       detail=revoked_detail)
        return Response({"ok": True, "grant_type": grant_type, "grant_id": grant_id})

    # ------------------------------------------------------------------
    # 访问申请单：我的申请 / 待我审批 / 审批
    # ------------------------------------------------------------------
    @action(detail=False, methods=["get"], url_path="my_access_requests")
    def my_access_requests(self, request):
        """GET /documents/my_access_requests/  我发起的访问申请

        通过统一工单（TicketList + TicketPermissionDetail）管理审批工单：
        文档目标信息从 reason 前缀 [doc:{id}:{action}] 解析。
        """
        qs = TicketList.objects.filter(
            applicant=request.user,
            permission_detail__change_type=TicketChangeType.GRANT,
        ).select_related('permission_detail').order_by("-created_at")[:100]
        data = []
        for ticket in qs:
            target_type, target_id, action, user_reason = _decode_ticket_reason(ticket.reason)
            data.append({
                "id": ticket.id,
                "target_type": target_type,
                "target_id": target_id,
                "action": action,
                "reason": user_reason,
                "status": ticket.status,
                "reviewer_comment": _extract_last_comment(ticket.approval_chain),
                "created_at": ticket.created_at,
                "updated_at": ticket.updated_at,
            })
        return Response(data)

    @action(detail=False, methods=["get"], url_path="pending_access_requests")
    def pending_access_requests(self, request):
        """GET /documents/pending_access_requests/  待我（所有者/管理员）审批的申请

        通过统一工单（TicketList + TicketPermissionDetail）管理审批工单。
        管理员看全部待审批工单；非管理员仅看自己文档对应的工单。
        """
        user = request.user
        is_manager = (getattr(user, 'is_super_admin', False)
                       or getattr(user, 'is_kb_admin', False))
        qs = TicketList.objects.filter(
            status=TicketStatus.PENDING,
            permission_detail__change_type=TicketChangeType.GRANT,
        ).select_related("applicant", "permission_detail")
        # 管理员看全部；非管理员仅看自己文档的申请（通过 reason 前缀匹配 doc_id）
        if not is_manager:
            owned_doc_ids = list(
                Document.objects.filter(owner=user, is_deleted=False).values_list('id', flat=True)
            )
            if owned_doc_ids:
                # 构造 reason 前缀匹配条件：[doc:{id}:...]
                import django.db.models as dm
                q = dm.Q()
                for did in owned_doc_ids:
                    q |= dm.Q(permission_detail__reason__startswith=f'[doc:{did}:')
                qs = qs.filter(q)
            else:
                qs = qs.none()
        qs = qs.order_by("-created_at")[:200]
        data = []
        for ticket in qs:
            target_type, target_id, action, user_reason = _decode_ticket_reason(ticket.reason)
            item = {
                "id": ticket.id,
                "applicant_id": ticket.applicant_id,
                "applicant_name": ticket.applicant.username if ticket.applicant else '',
                "target_type": target_type,
                "target_id": target_id,
                "action": action,
                "reason": user_reason,
                "status": ticket.status,
                "reviewer_comment": _extract_last_comment(ticket.approval_chain),
                "created_at": ticket.created_at,
                "updated_at": ticket.updated_at,
            }
            data.append(item)
        return Response(data)

    @action(detail=False, methods=["post"], url_path="approve_access_request")
    def approve_access_request(self, request):
        """POST /documents/approve_access_request/  {request_id, comment?}  批准并创建授权

        通过统一工单（TicketList + TicketPermissionDetail）管理审批工单：
        - 双层审批通过 approval_chain（JSON 数组）记录每步审批人/意见/时间
        - 每步审批写 TicketFlowLog（APPROVE），最终通过后写 EXECUTE
        - 复核通过后写入 ResourceShare（文档级个人共享）
        - visibility_change 类型工单通过后修改 doc.visibility_level
        """
        req_id = request.data.get("request_id")
        try:
            ticket = TicketList.objects.select_related(
                'applicant', 'permission_detail',
            ).get(id=req_id, status=TicketStatus.PENDING)
        except TicketList.DoesNotExist:
            raise Http404("申请不存在或已处理")

        # 从 reason 解析目标信息（doc 文档工单 / node 节点可见范围工单）
        target_type, target_id, action, user_reason = _decode_ticket_reason(ticket.reason)

        # 仅所有者/管理员可审批
        doc = None
        node = None
        if target_type == 'node':
            # 节点可见范围变更工单：审批人必须匹配审批链当前步骤的 approver_role
            # （团队级=部门经理，部门级=文档管理员/超管），兼容旧工单的节点所有者审批
            node = KnowledgeNode.objects.filter(id=target_id, is_deleted=False).first()
            if not node:
                raise Http404("目标节点不存在")
            chain = list(ticket.approval_chain or [])
            current_step = ticket.current_step or 0
            step_info = chain[current_step] if current_step < len(chain) else {}
            if not _can_approve_node_visibility(request.user, node, step_info):
                raise PermissionDenied("无权审批此申请")
        elif target_id:
            doc = Document.objects.filter(id=target_id, is_deleted=False).first()
            a = resolve_doc_access(request.user, doc) if doc else None
            if not a or not (a["is_owner"] or a["is_manager"]):
                raise PermissionDenied("无权审批此申请")
        else:
            if not (getattr(request.user, 'is_super_admin', False)
                    or getattr(request.user, 'is_kb_admin', False)):
                raise PermissionDenied("只有管理员可以审批此类申请")

        reviewer = request.user
        comment = (request.data.get("comment") or "")[:1000]

        # 双层审批逻辑：approval_chain 有多个 step 时需逐步审批
        chain = list(ticket.approval_chain or [])
        need_double = len(chain) > 1

        with transaction.atomic():
            if need_double:
                current_step = ticket.current_step or 0
                step_info = chain[current_step] if current_step < len(chain) else {}
                # 审核：不能与已审步骤为同一人（复核校验）
                for prev in chain[:current_step]:
                    if prev.get('approver_id') == reviewer.id:
                        raise PermissionDenied("双层审批不能由同一管理员完成，请另一位管理员审批")
                # 记录当前步骤审批结果
                step_info['approver_id'] = reviewer.id
                step_info['status'] = 'approved'
                step_info['comment'] = comment
                step_info['approved_at'] = timezone.now().isoformat()
                chain[current_step] = step_info

                if current_step + 1 < len(chain):
                    # 还有后续步骤，保持 pending
                    ticket.approval_chain = chain
                    ticket.current_step = current_step + 1
                    ticket.save(update_fields=["approval_chain", "current_step"])
                    TicketFlowLog.objects.create(
                        ticket=ticket, action='APPROVE', actor=reviewer,
                        comment=comment, step=current_step,
                    )
                    logger.info(f"[AccessRequest] step {current_step} approved id={ticket.id} by={reviewer.username} (pending next step)")
                    return Response({
                        "id": ticket.id,
                        "status": TicketStatus.PENDING,
                        "message": f"第 {current_step + 1} 审已通过，等待后续审批",
                    })
                # 所有步骤完成，最终通过
                ticket.approval_chain = chain
                ticket.status = TicketStatus.APPROVED
                ticket.approved_at = timezone.now()
                ticket.save(update_fields=["approval_chain", "status", "approved_at"])
                TicketFlowLog.objects.create(
                    ticket=ticket, action='APPROVE', actor=reviewer,
                    comment=comment, step=current_step,
                )
            else:
                # 普通单层审批
                if chain:
                    chain[0]['approver_id'] = reviewer.id
                    chain[0]['status'] = 'approved'
                    chain[0]['comment'] = comment
                    chain[0]['approved_at'] = timezone.now().isoformat()
                ticket.approval_chain = chain
                ticket.status = TicketStatus.APPROVED
                ticket.approved_at = timezone.now()
                ticket.save(update_fields=["approval_chain", "status", "approved_at"])
                TicketFlowLog.objects.create(
                    ticket=ticket, action='APPROVE', actor=reviewer,
                    comment=comment, step=0,
                )

            # 审批通过后执行授权写入
            if action == 'visibility_change' and (doc or node):
                # visibility_change 工单：修改文档/节点可见性层级
                # 从 user_reason 中解析目标 visibility_level（创建工单时编码在 reason 文本中）
                # 兜底使用 PUBLIC（扩大审批通常目标就是全局公开）
                new_level = VisibilityLevel.PUBLIC
                m = re.search(r'目标值:(\w+)', user_reason or '')
                if m:
                    # 显式编码的目标值（INHERIT 表示继承父级，写回 NULL）
                    new_level = m.group(1) if m.group(1) in VisibilityLevel.values else None
                else:
                    # 兼容旧格式：从申请文本中匹配枚举值，未命中兜底 PUBLIC
                    for level in VisibilityLevel.values:
                        if level in (user_reason or ''):
                            new_level = level
                            break
                if node:
                    # 节点可见范围变更：写回节点（NULL=继承父级，工单目标通常为具体三档值）
                    node.visibility_level = new_level
                    node.save(update_fields=['visibility_level', 'updated_at'])
                    _log_operation(request, 'node_visibility_change', node=node,
                                   detail={'ticket_id': ticket.id, 'applicant': ticket.applicant.username,
                                           'new_visibility_level': new_level})
                else:
                    doc.visibility_level = new_level
                    doc.save(update_fields=['visibility_level', 'updated_at'])
                    _log_operation(request, 'doc_visibility_change', document=doc,
                                   detail={'ticket_id': ticket.id, 'applicant': ticket.applicant.username,
                                           'new_visibility_level': new_level})
            elif target_id and doc:
                # 文档访问申请：创建 ResourceShare 个人级共享
                ResourceShare.objects.get_or_create(
                    resource_type=ResourceType.DOCUMENT,
                    resource_id=target_id,
                    share_scope_type=ShareScopeType.USER,
                    share_scope_id=ticket.applicant_id,
                    defaults={
                        'access_level': AccessLevel.READ,
                        'granted_by': reviewer,
                        'status': ShareStatus.ACTIVE,
                    },
                )
                if not doc.has_resource_share:
                    doc.has_resource_share = True
                    doc.save(update_fields=['has_resource_share'])
                _log_operation(request, 'doc_grant', document=doc,
                               detail={'ticket_id': ticket.id, 'applicant': ticket.applicant.username,
                                       'action': action, 'type': 'allow_user'})

            # 标记工单已执行
            ticket.status = TicketStatus.EXECUTED
            ticket.executed_at = timezone.now()
            ticket.save(update_fields=['status', 'executed_at'])
            TicketFlowLog.objects.create(ticket=ticket, action='EXECUTE', actor=reviewer)

        logger.info(f"[AccessRequest] approved id={ticket.id} applicant={ticket.applicant.username} target={target_type}:{target_id}")
        return Response({
            "id": ticket.id,
            "status": ticket.status,
            "applicant_id": ticket.applicant_id,
            "target_type": target_type,
            "target_id": target_id,
            "action": action,
        })

    @action(detail=False, methods=["post"], url_path="reject_access_request")
    def reject_access_request(self, request):
        """POST /documents/reject_access_request/  {request_id, comment?}  驳回

        通过统一工单（TicketList + TicketPermissionDetail）管理审批工单。
        驳回时在 approval_chain 当前步骤记录驳回意见并写 TicketFlowLog(REJECT)，
        工单状态置为 REJECTED。
        """
        req_id = request.data.get("request_id")
        try:
            ticket = TicketList.objects.select_related(
                'applicant', 'permission_detail',
            ).get(id=req_id, status=TicketStatus.PENDING)
        except TicketList.DoesNotExist:
            raise Http404("申请不存在或已处理")
        target_type, target_id, action, user_reason = _decode_ticket_reason(ticket.reason)
        doc = None
        node = None
        if target_type == 'node':
            # 节点可见范围变更工单：节点所有者或管理员可审批
            node = KnowledgeNode.objects.filter(id=target_id, is_deleted=False).first()
            if not (node and (node.owner_user_id == request.user.id
                              or getattr(request.user, 'is_super_admin', False)
                              or getattr(request.user, 'is_kb_admin', False))):
                raise PermissionDenied("无权审批此申请")
        elif target_id:
            doc = Document.objects.filter(id=target_id, is_deleted=False).first()
            a = resolve_doc_access(request.user, doc) if doc else None
            if not a or not (a["is_owner"] or a["is_manager"]):
                raise PermissionDenied("无权审批此申请")
        else:
            if not (getattr(request.user, 'is_super_admin', False)
                    or getattr(request.user, 'is_kb_admin', False)):
                raise PermissionDenied("只有管理员可以审批此类申请")

        comment = (request.data.get("comment") or "")[:1000]
        # 在 approval_chain 当前步骤记录驳回
        chain = list(ticket.approval_chain or [])
        current_step = ticket.current_step or 0
        with transaction.atomic():
            if current_step < len(chain):
                chain[current_step]['approver_id'] = request.user.id
                chain[current_step]['status'] = 'rejected'
                chain[current_step]['comment'] = comment
                chain[current_step]['approved_at'] = timezone.now().isoformat()

            ticket.approval_chain = chain
            ticket.status = TicketStatus.REJECTED
            ticket.save(update_fields=["approval_chain", "status"])
            TicketFlowLog.objects.create(
                ticket=ticket, action='REJECT', actor=request.user,
                comment=comment, step=current_step,
            )
        _log_operation(request, 'doc_grant_reject', document=doc, node=node,
                       detail={'ticket_id': ticket.id, 'applicant': ticket.applicant.username,
                               'action': action, 'target_type': target_type, 'target_id': target_id})
        return Response({
            "id": ticket.id,
            "status": ticket.status,
        })


class DocumentUploadView(APIView):
    """
    POST /api/v1/knowledge/documents/upload/
    multipart/form-data: file, node_id, [title], [visibility_level], [force_upload]
    - sha256 去重
    - 存到 MEDIA_ROOT/documents/{uuid}_{name} 或 OSS
    - 触发 Celery 异步解析

    visibility_level: TEAM_ONLY / DEPT_ONLY / PUBLIC
    （兼容旧版 visible_scope: team / dept / public）

    TODO: 需要优化上传逻辑,增加 root_type(领域) 参数让用户显式选择领域,
    当前领域由 node_id 所属的根节点隐式决定,用户无法在上传时选择目标领域。
    """
    permission_classes = [IsAuthenticated]
    parser_classes = [MultiPartParser, FormParser, JSONParser]

    def post(self, request):
        f = request.FILES.get("file")
        node_id = request.data.get("node_id")
        if not f or not node_id:
            return Response({"detail": "file / node_id 必填"}, status=400)

        try:
            node = KnowledgeNode.objects.get(id=node_id, is_deleted=False)
        except KnowledgeNode.DoesNotExist:
            return Response({"detail": "node 不存在"}, status=404)

        # 文档只能上传到文件夹（FOLDER）；ROOT/ORG 节点只作分支，不可直接挂文档
        if node.node_kind and node.node_kind != 'FOLDER':
            return Response({"detail": "文档只能上传到文件夹中，请选择文件夹节点"}, status=400)

        if not self._check_node_upload_permission(request.user, node):
            return Response({"detail": "无权限向该节点上传文档"}, status=403)

        ext = os.path.splitext(f.name)[1].lower()
        if ext not in ALLOWED_EXTENSIONS:
            return Response({"detail": f"不支持的文件类型: {ext}"}, status=400)

        # 验证文件真实类型（防止文件伪装）
        try:
            # 读取文件开头部分进行类型检测
            file_content = f.read(2048)
            f.seek(0)  # 重置文件指针
            detected_mime = magic.from_buffer(file_content, mime=True)
            # 根据扩展名验证MIME类型
            # 注意：纯文本类文件（代码/配置等）在不同系统的libmagic版本下，
            # 可能被检测为 text/plain 而非特定类型，因此增加 text/plain 作为备选
            ext_mime_map = {
                '.pdf': 'application/pdf',
                '.doc': 'application/msword',
                '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                '.md': ['text/markdown', 'text/plain'],
                '.markdown': ['text/markdown', 'text/plain'],
                '.txt': 'text/plain',
                '.rst': ['text/x-rst', 'text/plain'],
                # 电子表格
                '.csv': ['text/csv', 'application/csv', 'text/plain'],
                '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                '.xls': 'application/vnd.ms-excel',
                # 演示文稿
                '.ppt': 'application/vnd.ms-powerpoint',
                '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                # WPS Office 格式（MIME 类型与对应 MS 格式一致）
                '.wps': ['application/msword', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'],
                '.et': ['application/vnd.ms-excel', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'],
                '.dps': ['application/vnd.ms-powerpoint', 'application/vnd.openxmlformats-officedocument.presentationml.presentation'],
                # 代码类
                '.py': ['text/x-python', 'text/plain'],
                '.java': ['text/x-java-source', 'text/plain'],
                '.go': ['text/x-go', 'text/plain'],
                '.js': ['application/javascript', 'text/javascript', 'text/plain'],
                '.ts': ['text/typescript', 'application/typescript', 'text/plain'],
                '.jsx': ['application/javascript', 'text/javascript', 'text/plain'],
                '.tsx': ['text/typescript', 'application/typescript', 'text/plain'],
                '.c': ['text/x-c', 'text/plain'],
                '.cpp': ['text/x-c++', 'text/plain'],
                '.h': ['text/x-c', 'text/x-c++', 'text/plain'],
                '.rs': ['text/x-rust', 'text/plain'],
                '.yaml': ['text/yaml', 'text/x-yaml', 'text/plain'],
                '.yml': ['text/yaml', 'text/x-yaml', 'text/plain'],
                '.json': ['application/json', 'text/plain'],
                '.xml': ['application/xml', 'text/xml', 'text/plain'],
                '.toml': ['text/toml', 'text/plain'],
                '.ini': ['text/x-ini', 'text/plain'],
                '.conf': 'text/plain',
                '.cfg': 'text/plain',
                '.sh': ['text/x-shellscript', 'text/plain'],
                '.bat': ['application/x-bat', 'text/plain'],
                '.ps1': ['text/x-powershell', 'text/plain'],
                '.css': ['text/css', 'text/plain'],
            }
            expected_mime = ext_mime_map.get(ext)
            if expected_mime:
                if isinstance(expected_mime, list):
                    if detected_mime not in expected_mime:
                        return Response({"detail": f"文件类型不匹配：扩展名显示为 {ext}，但实际文件类型为 {detected_mime}"}, status=400)
                else:
                    if detected_mime != expected_mime:
                        return Response({"detail": f"文件类型不匹配：扩展名显示为 {ext}，但实际文件类型为 {detected_mime}"}, status=400)
        except Exception as e:
            logger.error(f"文件类型检测失败: {e}")
            return Response({"detail": "文件类型检测失败，请上传合法文件"}, status=400)

        if f.size > MAX_FILE_SIZE:
            return Response({"detail": f"文件大小超过限制（最大 {MAX_FILE_SIZE//(1024*1024)} MB）"}, status=400)

        # 可见性层级：兼容新版 visibility_level（TEAM_ONLY/DEPT_ONLY/PUBLIC）
        # 和旧版 visible_scope（team/dept/public），统一归一化为 visibility_level
        raw_visibility = request.data.get("visibility_level") or request.data.get("visible_scope")
        visibility_level = _normalize_visibility_level(raw_visibility)
        if raw_visibility and visibility_level is None:
            # 显式指定了可见范围但值非法 → 400（区别于未指定时的继承语义）
            return Response(
                {"detail": "visibility_level 必须是 TEAM_ONLY/DEPT_ONLY/PUBLIC 之一"},
                status=400,
            )
        if visibility_level is None:
            # 未显式指定可见范围：继承挂载文件夹的可见性（沿祖先链取最近非空，root 兜底 PUBLIC）
            visibility_level = _resolve_node_visibility(node)

        # 上传时选择是否允许下载/分享（默认只读：仅预览/对话检索）
        allow_download = request.data.get("allow_download") in ("true", "True", "1", True)
        allow_share = request.data.get("allow_share") in ("true", "True", "1", True)

        # 验证用户是否有权限设置指定的可见性层级
        is_valid, error_msg = _validate_visibility_level(request.user, visibility_level)
        if not is_valid:
            return Response({"detail": error_msg}, status=403)

        visibility_depts = request.data.getlist("visibility_depts", [])
        visibility_teams = request.data.getlist("visibility_teams", [])

        title = request.data.get("title") or f.name
        file_type = _detect_file_type(f.name)

        # 从节点祖先链推导归属 dept_id/team_id（组织 ID，非节点 ID）
        # Level 2 节点 ref_id = dept.id，Level 3 节点 ref_id = team.id
        dept_id = None
        team_id = None
        if node.node_level >= 2:
            ancestors = []
            current = node
            while current:
                ancestors.append(current)
                current = current.parent
            ancestors.reverse()
            for n in ancestors:
                if n.node_level == 2 and n.ref_id:
                    dept_id = n.ref_id
                elif n.node_level == 3 and n.ref_id:
                    team_id = n.ref_id

        # 归属约束：team_id 或 dept_id 至少一个非空
        # 若节点祖先链未推导出 dept_id，回退到上传者的主部门
        if not dept_id:
            dept_id = getattr(request.user, 'department_id', None)
        # 若未推导出 team_id，回退到上传者的所属团队
        if not team_id:
            team_id = getattr(request.user, 'team_id', None)

        # 归属约束校验：节点无组织祖先且上传者也无部门/团队归属时，
        # 非 PUBLIC 可见性既无部门/团队可挂靠，也会违反 doc_owner_scope_required 约束。
        # 不能静默降级为 PUBLIC（会造成越权公开），直接报错让用户改选节点或明确公开
        if not dept_id and not team_id and visibility_level != VisibilityLevel.PUBLIC:
            return Response({
                "detail": "当前节点无组织归属，且您未加入任何部门/团队，"
                          "仅支持上传为「全局公开」文档，请切换可见范围或选择其他节点"
            }, status=400)

        # 计算文件哈希 + 截取内容样本（单次读取流，避免二次 IO）
        h = hashlib.sha256()
        total = 0
        sample_raw = bytearray()
        for c in f.chunks():
            h.update(c)
            total += len(c)
            if len(sample_raw) < VERSION_SAMPLE_MAX_BYTES:
                sample_raw.extend(c[:VERSION_SAMPLE_MAX_BYTES - len(sample_raw)])
        file_hash = h.hexdigest()
        content_sample = _capture_content_sample(bytes(sample_raw), file_type)

        # 版本判定按「同组」进行：node + file_name + dept_id + team_id。
        # 不同部门/团队上传同名文档互不干扰（各自独立），
        # 避免团队 A 的文档被团队 B 的同名同内容上传误置非活跃。
        group_filter = dict(node=node, file_name=f.name[:256],
                            dept_id=dept_id, team_id=team_id, is_deleted=False)

        version_tag = request.data.get("version_tag", "").strip()

        if not version_tag:
            max_version = Document.objects.filter(
                **group_filter
            ).aggregate(models.Max('version'))['version__max'] or 0
            version_tag = f'v{max_version + 1}'
            version = max_version + 1
        else:
            existing_with_tag = Document.objects.filter(
                version_tag=version_tag, **group_filter
            ).first()
            if existing_with_tag:
                version = existing_with_tag.version
            else:
                max_version = Document.objects.filter(
                    **group_filter
                ).aggregate(models.Max('version'))['version__max'] or 0
                version = max_version + 1

        # 同版本标签去重：软删旧记录后重建（dedup 语义与旧版一致，仅收敛到同组范围）
        exist = Document.objects.filter(
            version_tag=version_tag, **group_filter
        ).first()

        # 组内其余文档（不含将被去重替换的 exist）→ 判定本次上传是「新版本」还是「独立文档」
        siblings = [s for s in Document.objects.filter(**group_filter)
                    if s.id != (exist.id if exist else -1)]
        is_version_upload = _is_version_upload(file_type, content_sample, siblings)

        file_path = None
        doc = None

        try:
            with transaction.atomic():
                file_path = self._save_file(f, node)
                if not file_path:
                    raise Exception("文件存储失败")

                if exist:
                    exist.is_deleted = True
                    exist.delete_time = timezone.now()
                    exist.save(update_fields=["is_deleted", "delete_time", "updated_at"])

                doc = Document.objects.create(
                    node=node,
                    title=title[:256],
                    file_name=f.name[:256],
                    file_type=file_type,
                    file_size=total,
                    file_hash=file_hash,
                    file_path=file_path,
                    mime_type=(f.content_type or "")[:64],
                    owner=request.user,
                    # node(FK) + dept_id + team_id 二选一非空
                    dept_id=dept_id,
                    team_id=team_id,
                    # visibility_level 控制可见范围
                    visibility_level=visibility_level,
                    allow_download=allow_download,
                    allow_share=allow_share,
                    root_type=node.root_type,
                    status="pending",
                    version=version,
                    version_tag=version_tag,
                    content_sample=content_sample,
                    is_active=True,
                )

                # 新版本：同组旧版本自动置非活跃，并同步检索向量表的活跃标志
                if is_version_upload:
                    # 事务内按同组条件重新查询（排除刚创建的新文档）后批量置非活跃，
                    # 避免基于事务外预加载的 siblings 在并发上传下造成误判
                    deactivated_qs = Document.objects.filter(**group_filter).exclude(id=doc.id)
                    deactivated_ids = list(deactivated_qs.values_list('id', flat=True))
                    deactivated_qs.update(is_active=False)
                    _sync_vectors_active(deactivated_ids, False)

                # visibility_teams：为指定团队创建跨团队共享
                if visibility_teams:
                    from apps.users.models import Team
                    team_ids_to_share = list(Team.objects.filter(
                        id__in=visibility_teams, is_deleted=False
                    ).values_list('id', flat=True))
                    for share_team_id in team_ids_to_share:
                        ResourceShare.objects.get_or_create(
                            resource_type=ResourceType.DOCUMENT,
                            resource_id=doc.id,
                            share_scope_type=ShareScopeType.TEAM,
                            share_scope_id=share_team_id,
                            defaults={
                                'access_level': AccessLevel.READ,
                                'granted_by': request.user,
                                'status': ShareStatus.ACTIVE,
                            },
                        )
                    # has_resource_share 标记是否存在资源共享
                    doc.has_resource_share = True
                    doc.save(update_fields=["has_resource_share"])

                _log_operation(request, 'doc_upload', document=doc, node=node,
                               detail={'file_name': f.name, 'file_size': total, 'file_hash': file_hash,
                                       'visibility_level': visibility_level, 'version_tag': version_tag,
                                       'visibility_teams': visibility_teams, 'visibility_depts': visibility_depts,
                                       'dept_id': dept_id, 'team_id': team_id,
                                       'is_version_upload': is_version_upload})
        except Exception as e:
            if file_path:
                try:
                    storage = get_document_storage()
                    storage.delete(file_path)
                except Exception:
                    logger.exception(f"Failed to clean up orphan file: {file_path}")
            return Response({"detail": str(e)[:200]}, status=500)

        celery_ok = True
        celery_error = ""
        try:
            from apps.knowledge.tasks import parse_document
            parse_document.delay(doc.id)
        except Exception as e:
            celery_ok = False
            celery_error = str(e)[:200]
            logger.warning(f"celery unreachable, doc {doc.id} will need manual reparse: {celery_error}")

        return Response({
            "document_id": doc.id,
            "uuid": str(doc.uuid),
            "status": doc.status,
            "file_hash": file_hash,
            "dedup": bool(exist),
            "celery_ok": celery_ok,
            "celery_error": celery_error,
            "version": doc.version,
            "is_active": doc.is_active,
            "is_version_upload": is_version_upload,
        }, status=201)

    def _check_node_upload_permission(self, user, node):
        """校验用户是否有权向目标节点上传文档

        超管/文档管理员：任意位置；
        部门经理：本部门节点（ORG）及其所有后代文件夹（path 前缀匹配）；
        团队组长/contributor：本团队节点（ORG）及其所有后代文件夹。
        """
        if getattr(user, 'is_super_admin', False) or getattr(user, 'is_kb_admin', False):
            return True

        role, _dept_id, team_ids = _get_user_role(user)

        if role == 'dept_manager':
            for dp in _get_dept_node_paths(user):
                if node.path == dp or node.path.startswith(dp):
                    return True

        if role in ('team_leader', 'contributor') and team_ids:
            for tp in _get_team_node_paths(team_ids):
                if node.path == tp or node.path.startswith(tp):
                    return True

        return False

    def _save_file(self, f, node):
        from apps.knowledge.storage import get_document_storage, generate_node_storage_path
        storage = get_document_storage()
        # 使用 Django 的 get_valid_filename 处理文件名，移除危险字符
        safe_name = django_text.get_valid_filename(f.name)
        # 进一步清理：移除控制字符和其他危险字符
        safe_name = re.sub(r'[\x00-\x1f\x7f]', '', safe_name)
        safe_name = safe_name.replace("..", "_")
        if not safe_name:
            safe_name = "unnamed_file"
        fname = f"{uuid_lib.uuid4().hex}_{safe_name}"
        # 生成节点存储路径
        node_path = generate_node_storage_path(node)
        logger.info(f'[Upload] saving file to node_path={node_path}, filename={fname}, node_id={node.id}, node_name={node.name}')
        file_path = storage.save(fname, f, node_path)
        logger.info(f'[Upload] file saved to: {file_path}')
        return file_path


class DocumentChunksView(APIView):
    """GET /api/v1/knowledge/documents/{id}/chunks/"""
    permission_classes = [IsAuthenticated]

    def get(self, request, doc_id):
        # 调试日志：确认用户认证状态
        logger.info(f'[Chunks] request user: {request.user}, is_authenticated: {request.user.is_authenticated}, is_super_admin: {getattr(request.user, "is_super_admin", False)}')
        
        try:
            doc = Document.objects.get(id=doc_id, is_deleted=False)
        except Document.DoesNotExist:
            return Response({"detail": "文档不存在"}, status=404)
        
        if not resolve_doc_access(request.user, doc)["can_read"]:
            logger.warning(f'[Chunks] user {request.user} has no permission for doc {doc_id}')
            return Response({"detail": "无权限查看此文档"}, status=403)

        chunks = DocumentChunk.objects.filter(document_id=doc_id).order_by("chunk_index")[:500]
        chunk_list = list(chunks)
        return Response({
            "document_id": int(doc_id),
            "total": len(chunk_list),
            "chunks": DocumentChunkSerializer(chunk_list, many=True).data,
        })


class CeleryStatusView(APIView):
    """GET /api/v1/knowledge/celery/status/ — 检查文档解析服务状态"""
    permission_classes = [IsAuthenticated]

    def get(self, request):
        from rag_project.celery import app
        from django.conf import settings
        import time
        import redis

        broker_url = getattr(settings, 'CELERY_BROKER_URL', '')
        result_backend = getattr(settings, 'CELERY_RESULT_BACKEND', '')

        diagnostics = {}

        try:
            conn = app.connection_for_read()
            conn.ensure_connection(max_retries=2)
            conn.close()
            diagnostics['broker_connected'] = True
        except Exception as e:
            logger.warning(f"celery broker connection failed: {e}")
            return Response({
                "celery_ok": False,
                "detail": "消息队列（Redis）连接失败",
            }, status=200)

        try:
            response = app.control.ping(timeout=5)
            if response:
                worker_count = len(response)
                return Response({
                    "celery_ok": True,
                    "detail": f"文档解析服务运行正常（{worker_count} 个在线）",
                    "worker_count": worker_count,
                })
        except Exception as e:
            logger.warning(f"celery control.ping failed: {e}")

        try:
            r = redis.Redis.from_url(broker_url)
            queues = ['default', 'parse', 'memory', 'email']
            queue_lengths = {}
            for q in queues:
                length = r.llen(q)
                queue_lengths[q] = int(length)

            if queue_lengths.get('default', 0) == 0 and queue_lengths.get('parse', 0) == 0:
                test_result = app.send_task('rag_project.celery.debug_task', args=[], queue='default')
                time.sleep(6)
                result_ready = test_result.ready()

                if result_ready:
                    return Response({
                        "celery_ok": True,
                        "detail": "文件解析服务运行正常",
                        "worker_count": 1,
                    })
        except Exception as e:
            logger.warning(f"celery redis check failed: {e}")

        return Response({
            "celery_ok": False,
            "detail": "消息队列连接正常，但文档解析服务未运行",
        }, status=200)


class PendingDocsView(APIView):
    """GET /api/v1/knowledge/documents/pending/ — 获取待处理文档列表"""
    permission_classes = [IsAuthenticated]

    def get(self, request):
        # 返回所有进行中的文档（不仅是 pending，还包括 parsing, embedding 等）
        processing_statuses = ["pending", "parsing", "desensitizing", "chunking", "embedding", "embedding_failed"]
        pending_query = Document.objects.filter(
            is_deleted=False,
            status__in=processing_statuses,
            owner=request.user
        ).order_by("-created_at")
        pending_count = pending_query.count()
        pending = pending_query[:20]
        serializer = DocumentSerializer(pending, many=True)
        return Response({
            "total": pending_count,
            "documents": serializer.data,
        })

    def post(self, request):
        """重新触发当前用户待处理文档的解析任务"""
        # 支持重新触发 pending 和 embedding_failed 的文档
        pending_query = Document.objects.filter(
            is_deleted=False,
            status__in=["pending", "embedding_failed"],
            owner=request.user
        )
        pending = list(pending_query)
        count = 0
        failed = []
        try:
            from apps.knowledge.tasks import parse_document
            for doc in pending:
                try:
                    parse_document.delay(doc.id)
                    count += 1
                except Exception as e:
                    failed.append({"doc_id": doc.id, "error": str(e)[:100]})
        except Exception as e:
            return Response({
                "ok": False,
                "detail": "Celery 连接失败",
                "error": str(e)[:200],
            }, status=500)
        return Response({
            "ok": True,
            "total_pending": len(pending),
            "retriggered": count,
            "failed": failed,
        })


# ============================================================================
# 文档审核（双审：团队组长审核 → 部门经理/合规复核）
# ============================================================================

class DocAuditPendingView(APIView):
    """GET /api/v1/knowledge/documents/pending-audits/
    获取待当前用户审核的文档列表

    审核规则（对齐 Document.AUDIT_STATUS_CHOICES）：
    - pending_team: 待团队组长审核 → 文档 team_id 对应当前用户为团队 leader
    - pending_compliance: 待合规复核 → 部门 leader / kb_admin / super_admin
    """
    permission_classes = [IsAuthenticated]

    def get(self, request):
        user = request.user
        # 访问入口权限校验：仅知识管理员 / 部门经理 / 团队组长 / 超管可访问
        from apps.users.models import Department, Team, has_permission
        if not (user.is_super_admin
                or user.is_kb_admin
                or has_permission(user, 'kb.manage_all')
                or has_permission(user, 'kb.manage')):
            is_leader = (
                Team.objects.filter(leader_id=user.id, is_deleted=False).exists()
                or Department.objects.filter(leader_id=user.id, is_deleted=False).exists()
            )
            if not is_leader:
                raise PermissionDenied("无文档审核权限")

        qs = Document.objects.filter(
            is_deleted=False,
            audit_status__in=['pending_team', 'pending_compliance'],
        ).select_related('owner', 'node').order_by('-created_at')

        # 按用户身份 + audit_status 过滤范围
        rows = []
        for doc in qs:
            can_audit = False
            audit_step = ''
            if doc.audit_status == 'pending_team':
                # 团队组长审核：文档 team_id 对应团队 leader
                if doc.team_id and Team.objects.filter(
                    id=doc.team_id, leader_id=user.id, is_deleted=False
                ).exists():
                    can_audit = True
                    audit_step = '审核（团队组长）'
                # 用户显式拥有 kb_admin / super_admin：也可以审（兜底，防止团队 leader 空缺）
                elif user.is_super_admin or user.is_kb_admin:
                    can_audit = True
                    audit_step = '审核（管理员代审）'
            elif doc.audit_status == 'pending_compliance':
                # 复核：部门经理 / kb_admin / super_admin
                if user.is_super_admin or user.is_kb_admin:
                    can_audit = True
                    audit_step = '复核（合规审核）'
                elif doc.dept_id and Department.objects.filter(
                    id=doc.dept_id, leader_id=user.id, is_deleted=False
                ).exists():
                    can_audit = True
                    audit_step = '复核（部门经理）'
            if not can_audit:
                continue

            # 解析节点路径（部门/团队）
            dept_name = ''
            team_name = ''
            if doc.team_id:
                t = Team.objects.filter(id=doc.team_id, is_deleted=False).only('name', 'department_id').first()
                if t:
                    team_name = t.name
                    d = Department.objects.filter(id=t.department_id, is_deleted=False).only('name').first()
                    if d:
                        dept_name = d.name
            elif doc.dept_id:
                d = Department.objects.filter(id=doc.dept_id, is_deleted=False).only('name').first()
                if d:
                    dept_name = d.name

            rows.append({
                'id': doc.id,
                'uuid': str(doc.uuid),
                'title': doc.title,
                'file_name': doc.file_name,
                'file_type': doc.file_type,
                'file_size': doc.file_size,
                'visibility_level': doc.visibility_level,
                'secret_level': doc.secret_level,
                'audit_status': doc.audit_status,
                'audit_step': audit_step,
                'owner_id': doc.owner_id,
                'owner_name': doc.owner.real_name or doc.owner.username if doc.owner else '',
                'owner_email': doc.owner.email if doc.owner else '',
                'node_id': doc.node_id,
                'node_name': doc.node.name if doc.node else '',
                'dept_name': dept_name,
                'team_name': team_name,
                'version': doc.version,
                'version_tag': doc.version_tag or '',
                'created_at': doc.created_at.isoformat() if doc.created_at else '',
                'updated_at': doc.updated_at.isoformat() if doc.updated_at else '',
            })

        return Response({
            'rows': rows,
            'count': len(rows),
        })


class DocAuditApproveView(APIView):
    """POST /api/v1/knowledge/documents/<id>/audit-approve/
    审核通过文档

    状态流转：
    - pending_team → pending_compliance（审核通过，进入复核）
    - pending_compliance → passed（复核通过，正式放行）
    """
    permission_classes = [IsAuthenticated]

    @transaction.atomic
    def post(self, request, pk):
        from apps.users.models import Department, Team, has_permission

        comment = (request.data.get("comment") or "").strip()
        try:
            doc = Document.objects.select_for_update().get(pk=pk, is_deleted=False)
        except Document.DoesNotExist:
            return Response({"detail": "文档不存在"}, status=404)

        user = request.user
        # 权限 + 状态 校验
        if doc.audit_status == 'pending_team':
            can = False
            if doc.team_id and Team.objects.filter(
                id=doc.team_id, leader_id=user.id, is_deleted=False
            ).exists():
                can = True
            elif user.is_super_admin or user.is_kb_admin:
                can = True
            if not can:
                raise PermissionDenied("您不是该文档所属团队的组长，无权进行审核")
            # 审核通过 → 进入复核
            doc.audit_status = 'pending_compliance'
        elif doc.audit_status == 'pending_compliance':
            can = False
            if user.is_super_admin or user.is_kb_admin:
                can = True
            elif doc.dept_id and Department.objects.filter(
                id=doc.dept_id, leader_id=user.id, is_deleted=False
            ).exists():
                can = True
            if not can:
                raise PermissionDenied("您没有该文档的复核权限")
            # 复核通过 → passed
            doc.audit_status = 'passed'
        else:
            return Response({"detail": f"文档当前状态 {doc.audit_status} 不可审核"}, status=400)

        doc.save(update_fields=['audit_status', 'updated_at'])
        _log_operation(request, f'doc_audit_approve_{doc.audit_status}', document=doc,
                       detail={'comment': comment, 'approver': user.username})
        logger.info(f"Doc audit approved: id={pk}, status={doc.audit_status}, approver={user.username}")
        return Response({
            "ok": True,
            "audit_status": doc.audit_status,
            "title": doc.title,
        })


class DocAuditRejectView(APIView):
    """POST /api/v1/knowledge/documents/<id>/audit-reject/
    审核驳回文档（驳回理由必填）

    状态流转：pending_team / pending_compliance → rejected
    """
    permission_classes = [IsAuthenticated]

    @transaction.atomic
    def post(self, request, pk):
        from apps.users.models import Department, Team

        comment = (request.data.get("comment") or "").strip()
        if not comment:
            return Response({"detail": "驳回理由不能为空"}, status=400)

        try:
            doc = Document.objects.select_for_update().get(pk=pk, is_deleted=False)
        except Document.DoesNotExist:
            return Response({"detail": "文档不存在"}, status=404)

        user = request.user
        # 权限校验：与 approve 相同
        if doc.audit_status not in ('pending_team', 'pending_compliance'):
            return Response({"detail": f"文档当前状态 {doc.audit_status} 不可驳回"}, status=400)

        can = False
        if doc.audit_status == 'pending_team':
            if doc.team_id and Team.objects.filter(
                id=doc.team_id, leader_id=user.id, is_deleted=False
            ).exists():
                can = True
            elif user.is_super_admin or user.is_kb_admin:
                can = True
        else:  # pending_compliance
            if user.is_super_admin or user.is_kb_admin:
                can = True
            elif doc.dept_id and Department.objects.filter(
                id=doc.dept_id, leader_id=user.id, is_deleted=False
            ).exists():
                can = True
        if not can:
            raise PermissionDenied("您没有该文档的审核权限")

        doc.audit_status = 'rejected'
        doc.save(update_fields=['audit_status', 'updated_at'])
        _log_operation(request, 'doc_audit_reject', document=doc,
                       detail={'comment': comment, 'rejector': user.username})
        logger.info(f"Doc audit rejected: id={pk}, rejector={user.username}, reason={comment[:100]}")
        return Response({
            "ok": True,
            "audit_status": doc.audit_status,
            "title": doc.title,
            "reject_comment": comment,
        })
